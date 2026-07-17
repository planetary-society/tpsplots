"""Base chart generation view component with desktop/mobile versions built in."""

import csv
import logging
import math
import struct
import textwrap
import warnings
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import ClassVar

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.ticker import FuncFormatter

from tpsplots import TPS_STYLE_FILE
from tpsplots.colors import COLORS, TPS_COLORS, resolve_color
from tpsplots.exceptions import RenderingError
from tpsplots.views.mixins import AxisTickFormatMixin
from tpsplots.views.style import tokens

logger = logging.getLogger(__name__)


class ChartView(AxisTickFormatMixin):
    """Base class for all chart views with shared functionality."""

    # Pydantic config model for this chart type (set by subclasses)
    CONFIG_CLASS: ClassVar[type | None] = None

    # Color palettes imported from tpsplots.colors for backward compatibility
    COLORS: ClassVar[dict[str, str]] = COLORS
    TPS_COLORS: ClassVar[dict[str, str]] = TPS_COLORS
    # Device-specific visual settings
    DESKTOP: ClassVar[dict[str, object]] = {
        "type": "desktop",
        "figsize": (16, 10),
        "dpi": 300,
        "title_size": 28,
        "label_size": 16,
        "tick_size": 20,
        "legend_size": 18,
        "line_width": 3,
        "marker_size": 6,
        "grid": True,
        "grid_axis": "y",
        "tick_rotation": 0,
        "add_logo": True,
        "max_ticks": 25,
        "decade_tick_threshold": 50,  # Year ranges above this show decade labels only
        "minor_tick_length": 0,
        "minor_tick_width": 0,
        "major_tick_length": 4,
        "major_tick_width": 1,
        "footer": True,
        "footer_height": 0.08,
        "footer_line_width": 1,
        "footer_extent": (0.01, 0.99),
        "header": True,
        "header_height": 0.1,
        "header_min_height": 0.06,  # Never smaller than 6%
        "header_max_height": 0.18,  # Never larger than 18%
        "header_x": 0.01,
        "header_y": 0.98,
        "header_padding": 0.03,
        "chart_vertical_padding": 0.01,
        "subtitle_size_ratio": 0.7,
        "subtitle_line_spacing": 1.05,
        "title_line_spacing": 1.1,
        "subtitle_vertical_padding_scale": 0.5,
        "title_subtitle_gap": 0.005,
        "show_eyebrow": True,
        "subtitle_wrap_length": 120,
        "label_wrap_length": 30,
        "logo_zoom": 0.03,
        "logo_x": 0.009,
        "logo_y": 0.005,
        "source_size": 12,
        "source_y": 0.01,
    }

    MOBILE: ClassVar[dict[str, object]] = {
        "type": "mobile",
        "figsize": (8, 9),
        "dpi": 300,
        "title_size": 24,
        "label_size": 11,
        "tick_size": 13,
        "legend_size": 15,
        "line_width": 3,
        "marker_size": 5,
        "grid": True,
        "grid_axis": "y",
        "tick_rotation": 0,
        "add_logo": True,
        "decade_tick_threshold": 50,
        "minor_tick_length": 0,
        "minor_tick_width": 0,
        "major_tick_length": 4,
        "major_tick_width": 1,
        "footer": True,
        "footer_height": 0.08,
        "footer_line_width": 1,
        "footer_extent": (0.01, 0.99),
        "header": True,
        "header_height": 0.14,
        "header_min_height": 0.08,  # Never smaller than 8% (mobile needs more)
        "header_max_height": 0.22,  # Never larger than 22%
        "header_x": 0.01,
        "header_y": 0.98,
        "header_padding": 0.03,
        "chart_vertical_padding": 0.01,
        "subtitle_size_ratio": 0.7,
        "subtitle_line_spacing": 1.05,
        "title_line_spacing": 1.1,
        "subtitle_vertical_padding_scale": 0.5,
        "title_subtitle_gap": 0.005,
        "show_eyebrow": True,
        "subtitle_wrap_length": 64,
        "label_wrap_length": 15,
        "logo_zoom": 0.025,
        "logo_x": 0.008,
        "logo_y": 0.0065,
        "source_size": 9,
        "source_y": 0.0105,
    }

    SOCIAL: ClassVar[dict[str, object]] = {
        "type": "social",
        "figsize": (8, 4.2),  # 8*300=2400px, 4.2*300=1260px (40:21 OG ratio)
        "dpi": 300,  # match desktop/mobile; 1200x630 at native 150dpi looked soft
        "title_size": 18,
        "label_size": 9,
        "tick_size": 10,
        "legend_size": 11,
        "line_width": 3,
        "marker_size": 5,
        "grid": True,
        "grid_axis": "y",
        "tick_rotation": 0,
        "add_logo": True,
        "max_ticks": 15,
        "decade_tick_threshold": 50,
        "minor_tick_length": 0,
        "minor_tick_width": 0,
        "major_tick_length": 4,
        "major_tick_width": 1,
        "footer": True,
        "footer_height": 0.12,  # Larger than desktop/mobile to compensate for no header
        "footer_line_width": 1,
        "footer_extent": (0.01, 0.99),
        "header": False,
        "header_height": 0,
        "header_min_height": 0,
        "header_max_height": 0,
        "header_x": 0.01,
        "header_y": 0.98,
        "header_padding": 0.03,
        "chart_vertical_padding": 0.01,
        "subtitle_size_ratio": 0.7,
        "subtitle_line_spacing": 1.05,
        "title_line_spacing": 1.1,
        "subtitle_vertical_padding_scale": 0.5,
        "title_subtitle_gap": 0.005,
        "show_eyebrow": False,
        "subtitle_wrap_length": 80,
        "label_wrap_length": 25,
        "logo_zoom": 0.05,
        "logo_x": 0.009,
        "logo_y": 0.005,
        "source_size": 9,
        "source_y": 0.01,
    }

    # Video style variants render the chart panel ONLY (no header/title/subtitle,
    # no footer/logo/source) — titles and branding are composited in the video
    # editor, which is what lets the chart fill a 1:1 / 16:9 / 9:16 frame edge to
    # edge. header/footer/add_logo are all False (same mechanism as SOCIAL, which
    # already exercises header: False). Key parity with DESKTOP is required (views
    # read many of these keys) so header_*/title_size etc. are kept even though the
    # chrome they configure is never drawn. Consumed by ``tpsplots animate``.
    # The three video variants are ONE style at three aspect ratios, so shared
    # values live in _VIDEO_BASE and each variant lists only its real deltas.
    _VIDEO_BASE: ClassVar[dict[str, object]] = {
        "dpi": 150,
        "title_size": 20,
        "label_size": 16,
        "legend_size": 15,
        "marker_size": 9,
        "grid": True,
        "grid_axis": "y",
        "tick_rotation": 0,
        "add_logo": False,
        "decade_tick_threshold": 50,
        "minor_tick_length": 0,  # sub-pixel minors become yuv420p chroma mush
        "minor_tick_width": 0,
        "major_tick_length": 4,
        "major_tick_width": 1,
        "footer": False,
        "footer_height": 0,
        "footer_line_width": 1,
        "footer_extent": (0.01, 0.99),
        "header": False,
        "header_height": 0,
        "header_min_height": 0,
        "header_max_height": 0,
        "header_x": 0.01,
        "header_y": 0.98,
        "header_padding": 0.03,
        "chart_vertical_padding": 0.01,
        "subtitle_size_ratio": 0.7,
        "subtitle_line_spacing": 1.05,
        "title_line_spacing": 1.1,
        "subtitle_vertical_padding_scale": 0.5,
        "title_subtitle_gap": 0.005,
        "show_eyebrow": False,
        "logo_zoom": 0.05,
        "logo_x": 0.009,
        "logo_y": 0.005,
        "source_size": 10,
        "source_y": 0.01,
    }

    VIDEO_SQUARE: ClassVar[dict[str, object]] = {
        **_VIDEO_BASE,
        "type": "video_square",
        "figsize": (7.2, 7.2),  # 7.2*150=1080px, 7.2*150=1080px
        "tick_size": 15,
        "line_width": 4.5,
        "max_ticks": 8,
        "subtitle_wrap_length": 60,
        "label_wrap_length": 18,
    }

    VIDEO_LANDSCAPE: ClassVar[dict[str, object]] = {
        **_VIDEO_BASE,
        "type": "video_landscape",
        "figsize": (12.8, 7.2),  # 12.8*150=1920px, 7.2*150=1080px
        "tick_size": 16,
        "line_width": 5,
        "max_ticks": 12,
        "subtitle_wrap_length": 96,
        "label_wrap_length": 30,
    }

    VIDEO_PORTRAIT: ClassVar[dict[str, object]] = {
        **_VIDEO_BASE,
        "type": "video_portrait",
        "figsize": (7.2, 12.8),  # 7.2*150=1080px, 12.8*150=1920px
        "tick_size": 15,
        "line_width": 4.5,
        "max_ticks": 8,
        "subtitle_wrap_length": 60,
        "label_wrap_length": 14,
    }

    # Device name -> style-dict attribute, resolved via getattr so subclass
    # overrides of DESKTOP/MOBILE/etc. keep working. Unknown devices fall back
    # to DESKTOP (existing behavior).
    _DEVICE_STYLES: ClassVar[dict[str, str]] = {
        "desktop": "DESKTOP",
        "mobile": "MOBILE",
        "social": "SOCIAL",
        "video_square": "VIDEO_SQUARE",
        "video_landscape": "VIDEO_LANDSCAPE",
        "video_portrait": "VIDEO_PORTRAIT",
    }

    def __init__(self, outdir: Path = Path("charts"), style_file=TPS_STYLE_FILE):
        """
        Initialize the chart view with output directory and style.

        Args:
            outdir: Output directory for chart files
            style_file: Matplotlib style file path to use
        """
        self.outdir = outdir
        self.outdir.mkdir(parents=True, exist_ok=True)

        # Apply style if provided
        if style_file:
            plt.style.use(style_file)

    def device_style(self, device: str) -> dict:
        """The style dict for ``device``, raising on unknown names.

        The strict counterpart to ``create_figure``'s silent DESKTOP fallback —
        callers that must guarantee exact output dimensions (the animate
        renderer) use this instead of reaching into ``_DEVICE_STYLES``.
        """
        attr = self._DEVICE_STYLES.get(device)
        if attr is None:
            valid = ", ".join(sorted(self._DEVICE_STYLES))
            raise ValueError(f"Unknown device {device!r}. Valid devices: {valid}.")
        return getattr(self, attr)

    def create_figure(self, metadata, device="desktop", **kwargs):
        """Create a single chart figure for the given device.

        Unlike ``generate_chart``, this does not save files or create
        all device variants — it returns one matplotlib Figure for
        in-memory use (e.g. editor previews).

        Args:
            metadata: Chart metadata dictionary
            device: ``"desktop"``, ``"mobile"``, ``"social"``,
                ``"video_square"``, ``"video_landscape"``, or ``"video_portrait"``
            **kwargs: Additional parameters for chart creation

        Returns:
            matplotlib.figure.Figure: The created figure
        """
        kwargs.pop("export_data", None)
        style = getattr(self, self._DEVICE_STYLES.get(device, "DESKTOP"))
        chart_kwargs = self._clone_chart_kwargs(kwargs)
        chart_kwargs["style"] = style
        return self._create_chart_with_overlays(metadata, **chart_kwargs)

    def generate_chart(self, metadata, stem, **kwargs):
        """
        Generate desktop, mobile, and social versions of a chart.

        Args:
            metadata: Chart metadata dictionary
            stem: Base filename for the chart
            **kwargs: Additional parameters for chart creation

        Returns:
            dict: ``{"files": [...]}`` — paths of every generated output file.
            The per-device figures are saved and closed internally.
        """

        export_data = kwargs.pop("export_data", None)
        generated_files: list[str] = []

        try:
            # Create desktop version
            desktop_kwargs = self._clone_chart_kwargs(kwargs)
            desktop_kwargs["style"] = self.DESKTOP
            desktop_fig = self._create_chart_with_overlays(metadata, **desktop_kwargs)
            generated_files.extend(
                self._save_chart(desktop_fig, f"{stem}_desktop", metadata, create_pptx=True)
            )

            # Create mobile version
            mobile_kwargs = self._clone_chart_kwargs(kwargs)
            mobile_kwargs["style"] = self.MOBILE
            mobile_fig = self._create_chart_with_overlays(metadata, **mobile_kwargs)
            generated_files.extend(
                self._save_chart(mobile_fig, f"{stem}_mobile", metadata, create_pptx=False)
            )

            # Create social card version (PNG only, no header/footer)
            social_kwargs = self._clone_chart_kwargs(kwargs)
            social_kwargs["style"] = self.SOCIAL
            social_fig = self._create_chart_with_overlays(metadata, **social_kwargs)
            generated_files.extend(
                self._save_chart(social_fig, f"{stem}_social", metadata, create_svg=False)
            )

            # Export CSV if export_data is present
            if export_data is not None:
                csv_path = self._export_csv(export_data, metadata, stem)
                generated_files.append(str(csv_path))

        except Exception as e:  # Boundary: convert any render failure to RenderingError
            logger.error(f"Error generating chart {stem}: {e}")
            raise RenderingError(f"Error generating chart {stem}: {e}") from e

        logger.info(f"✓ generated charts for {stem}")

        # The desktop/mobile/social figures are already saved and closed by
        # ``_save_chart``; no caller consumes them, so only report the files.
        return {"files": generated_files}

    @staticmethod
    def _clone_chart_kwargs(kwargs: dict) -> dict:
        """Clone nested container kwargs so desktop/mobile renders cannot mutate shared state."""
        return {
            key: deepcopy(value)
            if isinstance(value, (pd.DataFrame, pd.Series, dict, list, tuple, set))
            else value
            for key, value in kwargs.items()
        }

    def _create_chart(self, metadata, style, **kwargs):
        """
        Abstract method to create a chart with the specified style.

        Args:
            metadata: Chart metadata dictionary
            style: Style dictionary (DESKTOP or MOBILE)
            **kwargs: Additional parameters for chart creation

        Returns:
            matplotlib.figure.Figure: The created figure

        Raises:
            NotImplementedError: This method must be implemented by subclasses
        """
        raise NotImplementedError("Subclasses must implement _create_chart")

    def _create_chart_with_overlays(self, metadata, **chart_kwargs):
        """Render a chart and apply figure-level overlays that every device shares.

        The single choke point wrapping ``_create_chart`` so that data-space
        annotations are drawn identically for desktop/mobile/social/video and
        editor previews, right after the subclass builds the figure.
        """
        fig = self._create_chart(metadata, **chart_kwargs)
        self._apply_annotations(fig, metadata, chart_kwargs.get("style"))
        return fig

    @staticmethod
    def _coerce_annotation_x(value):
        """Coerce a string annotation x to a Timestamp for date axes, else pass through.

        Charts with date x-axes need a datetime anchor; numeric axes keep the raw
        value. Non-string values (floats) are returned unchanged.
        """
        if isinstance(value, str):
            try:
                return pd.to_datetime(value)
            except (ValueError, TypeError):
                return value
        return value

    @staticmethod
    def _annotation_unit_x(ax, value):
        """Map a (already date-coerced) annotation x to the axis' numeric units.

        Runs the x through the axis unit converter so date Timestamps become
        matplotlib date numbers and categorical strings become their tick index,
        which is what ``transData`` expects. Unknown values pass through.
        """
        try:
            return ax.xaxis.convert_units(value)
        except Exception:  # defensive: fall back to the raw value
            return value

    def _apply_annotations(self, fig, metadata, style):
        """Draw editorial data-space callouts on the primary axes (``fig.axes[0]``).

        No-op unless ``metadata['annotations']`` is set. Each callout is text in
        the same white rounded box used by direct line labels
        (``tokens.direct_label_bbox``), optionally joined to its anchor by a thin
        curved ``drawarrow`` arrow.

        CONTRACT: arrows and flexitext frames are pinned in absolute display
        pixels (``IdentityTransform``), which is only correct because (a) all
        axes layout (`_adjust_layout_for_header_footer` and friends) is final
        before this runs, and (b) `_save_chart` saves with ``dpi="figure"`` and
        never ``bbox_inches="tight"``. Re-layout or dpi-rescaling after this
        method would silently detach every callout from its anchor.

        Rendering choices
        -----------------
        - **Text box.** Plain strings are drawn as a matplotlib ``Text`` whose
          built-in ``bbox`` gives the rounded frame. Strings carrying flexitext
          style tags (``<weight:semibold>...</>``) are drawn with ``flexitext``
          (which cannot itself draw a frame) and get a ``FancyBboxPatch`` sized to
          the laid-out text; ``mutation_scale = fontsize * dpi / 72`` makes that
          patch pixel-identical to a ``Text`` bbox, so both paths look the same.
        - **Placement.** Callouts with explicit ``text_x``/``text_y`` are honoured
          verbatim. The rest start at their anchor and are separated with
          ``adjustText`` so they don't overlap each other or the anchor points.
        - **Arrows.** Drawn last (after adjustment) in display coordinates so the
          gentle arc is uniform regardless of the axis scales, from the box edge
          to a small gap before the anchor.

        Args:
            fig: The matplotlib Figure object
            metadata: Chart metadata dictionary
            style: Style dictionary (for label font size); may be ``None``
        """
        annotations = metadata.get("annotations")
        if not annotations:
            return

        axes = fig.axes
        if not axes:
            return
        ax = axes[0]

        # Optional, relatively heavy deps: imported lazily so the (common)
        # no-annotation render never pays for them.
        from adjustText import adjust_text
        from drawarrow import ax_arrow
        from flexitext import flexitext
        from matplotlib.patches import BoxStyle, FancyBboxPatch
        from matplotlib.transforms import Bbox, IdentityTransform

        fontsize = (style or {}).get("label_size", 12)
        # boxstyle pad (0.2) is applied in display px scaled by mutation_scale;
        # replicate matplotlib's Text-bbox padding for the flexitext patches.
        mutation_scale = fontsize * fig.dpi / 72.0
        pad_frac = BoxStyle(tokens.DIRECT_LABEL_BBOX_BOXSTYLE).pad
        pad_px = pad_frac * mutation_scale

        # adjustText and the display-space patches must not rescale the view.
        xlim0, ylim0 = ax.get_xlim(), ax.get_ylim()

        records = []
        for raw in annotations:
            text = self._escape_svg_text(raw.get("text", "")) or ""
            text_x = raw.get("text_x")
            text_y = raw.get("text_y")
            explicit = text_x is not None and text_y is not None
            x = self._coerce_annotation_x(raw.get("x"))
            y = raw.get("y")
            pos = (self._coerce_annotation_x(text_x), text_y) if explicit else (x, y)
            records.append(
                {
                    "text": text,
                    "edge": resolve_color(raw.get("color")) or tokens.ANNOTATION_EDGE_COLOR,
                    "x": x,
                    "y": y,
                    "pos": pos,
                    "explicit": explicit,
                    "tagged": "</>" in text,
                    "arrow": bool(raw.get("arrow", False)),
                    "artist": None,
                    "is_flexi": False,
                    "box_bbox": None,
                }
            )

        # 1) Draw the text. Plain strings carry their rounded box directly; tagged
        #    strings get a matching FancyBboxPatch once flexitext has laid out.
        for rec in records:
            px, py = rec["pos"]
            if rec["tagged"]:
                styled = f"<color:{tokens.ANNOTATION_COLOR}, size:{fontsize}>{rec['text']}</>"
                try:
                    ab = flexitext(
                        px,
                        py,
                        styled,
                        xycoords="data",
                        ax=ax,
                        ha="center",
                        va="center",
                        ma="center",
                    )
                    ab.set_zorder(tokens.ANNOTATION_TEXT_ZORDER)
                    rec["artist"], rec["is_flexi"] = ab, True
                except Exception:  # malformed tags: fall back to plain text
                    rec["tagged"] = False
            if not rec["is_flexi"]:
                rec["artist"] = ax.text(
                    px,
                    py,
                    rec["text"],
                    fontsize=fontsize,
                    color=tokens.ANNOTATION_COLOR,
                    ha="center",
                    va="center",
                    multialignment="center",
                    zorder=tokens.ANNOTATION_TEXT_ZORDER,
                    bbox=tokens.direct_label_bbox(rec["edge"]),
                )

        # 2) Separate the label-less callouts; explicit coordinates stay put.
        movable = [r for r in records if not r["explicit"] and not r["is_flexi"]]
        if movable:
            static_objs = [r["artist"] for r in records if r not in movable]
            anchor_x = [self._annotation_unit_x(ax, r["x"]) for r in records]
            anchor_y = [ax.yaxis.convert_units(r["y"]) for r in records]
            adjust_text(
                [r["artist"] for r in movable],
                x=anchor_x,
                y=anchor_y,
                objects=static_objs or None,
                ax=ax,
                expand=(1.15, 1.3),
                force_text=(0.4, 0.6),
                only_move={"text": "xy", "static": "xy", "explode": "xy", "pull": "xy"},
                time_lim=0.4,
            )

        # 3) One draw so bbox patches exist and every extent is final/exact.
        fig.canvas.draw()
        renderer = fig.canvas.get_renderer()

        # 4) Boxes + final box extents (used as arrow tails).
        for rec in records:
            if rec["is_flexi"]:
                ext = rec["artist"].get_window_extent(renderer)
                patch = FancyBboxPatch(
                    (ext.x0, ext.y0),
                    ext.width,
                    ext.height,
                    mutation_scale=mutation_scale,
                    transform=IdentityTransform(),
                    clip_on=False,
                    zorder=tokens.ANNOTATION_BOX_ZORDER,
                    **tokens.direct_label_bbox(rec["edge"]),
                )
                ax.add_artist(patch)
                rec["box_bbox"] = Bbox.from_extents(
                    ext.x0 - pad_px, ext.y0 - pad_px, ext.x1 + pad_px, ext.y1 + pad_px
                )
            else:
                box_patch = rec["artist"].get_bbox_patch()
                rec["box_bbox"] = (
                    box_patch.get_window_extent(renderer)
                    if box_patch is not None
                    else rec["artist"].get_window_extent(renderer)
                )

        # 5) Arrows last, so they land on the FINAL box positions.
        for rec in records:
            if not rec["arrow"] or rec["box_bbox"] is None:
                continue
            try:
                x_num = self._annotation_unit_x(ax, rec["x"])
                y_num = ax.yaxis.convert_units(rec["y"])
                anchor = ax.transData.transform((x_num, y_num))
            except Exception:  # non-numeric anchor: skip the arrow
                continue
            self._draw_annotation_arrow(ax, rec["box_bbox"], anchor, rec["edge"], ax_arrow)

        ax.set_xlim(xlim0)
        ax.set_ylim(ylim0)

    @staticmethod
    def _draw_annotation_arrow(ax, box_bbox, anchor, edge, ax_arrow):
        """Draw a thin curved arrow from a callout box edge to its anchor point.

        All geometry is in display pixels (so the arc reads the same regardless
        of the axis scales); the arrow is then pinned to that display frame via
        ``IdentityTransform``. Returns the patch, or ``None`` when the box sits so
        close to its anchor that the shaft would be swallowed by the arrowhead.
        """
        from matplotlib.transforms import IdentityTransform

        # Point-based geometry scaled to the figure's pixel grid, so gaps read the
        # same at any dpi. The minimum shaft keeps a couple of head-lengths of line
        # beyond the arrowhead; below that a box sitting on its target just reads
        # as a bare triangle, so we skip the arrow and let proximity do the work.
        scale = ax.figure.dpi / 72.0
        tail_gap = 1.5 * scale
        head_gap = 3.0 * scale
        head_len_px = tokens.ANNOTATION_ARROW_HEAD_LENGTH * scale
        min_len = tail_gap + 2.2 * head_len_px

        cx = 0.5 * (box_bbox.x0 + box_bbox.x1)
        cy = 0.5 * (box_bbox.y0 + box_bbox.y1)
        dx, dy = anchor[0] - cx, anchor[1] - cy
        dist = math.hypot(dx, dy)
        if dist < 1e-6:
            return None
        ux, uy = dx / dist, dy / dist

        half_w = 0.5 * (box_bbox.x1 - box_bbox.x0)
        half_h = 0.5 * (box_bbox.y1 - box_bbox.y0)
        sx = half_w / abs(dx) if abs(dx) > 1e-9 else math.inf
        sy = half_h / abs(dy) if abs(dy) > 1e-9 else math.inf
        s = min(sx, sy)  # ray/box-boundary intersection parameter along (dx, dy)

        tail = (cx + s * dx + ux * tail_gap, cy + s * dy + uy * tail_gap)
        head = (anchor[0] - ux * head_gap, anchor[1] - uy * head_gap)
        if math.hypot(head[0] - tail[0], head[1] - tail[1]) < min_len:
            return None

        arrow = ax_arrow(
            list(tail),
            list(head),
            radius=tokens.ANNOTATION_ARROW_RADIUS,
            width=tokens.ANNOTATION_ARROW_WIDTH,
            head_width=tokens.ANNOTATION_ARROW_HEAD_WIDTH,
            head_length=tokens.ANNOTATION_ARROW_HEAD_LENGTH,
            ax=ax,
            color=edge,
            zorder=tokens.ANNOTATION_ARROW_ZORDER,
        )
        arrow.set_transform(IdentityTransform())
        arrow.set_clip_on(False)
        return arrow

    def _extract_metadata_from_kwargs(self, metadata: dict, kwargs: dict) -> None:
        """
        Extract title and subtitle from kwargs into metadata dict.

        This pattern is used by all chart views to allow title/subtitle to be
        passed either via metadata dict or as direct kwargs.

        Args:
            metadata: Chart metadata dictionary to update (modified in place)
            kwargs: Keyword arguments dict to extract from (modified in place)
        """
        for text in ["title", "subtitle"]:
            if kwargs.get(text):
                metadata[text] = kwargs.pop(text)

    def _setup_figure(self, style: dict, kwargs: dict) -> tuple:
        """
        Create figure and axes with style-appropriate sizing.

        This pattern is used by all chart views to create the figure with
        consistent figsize and dpi handling.

        Args:
            style: Style dictionary (DESKTOP or MOBILE)
            kwargs: Keyword arguments dict (figsize/dpi are popped if present)

        Returns:
            tuple: (fig, ax) matplotlib figure and axes objects
        """
        figsize = kwargs.pop("figsize", style["figsize"])
        dpi = kwargs.pop("dpi", style["dpi"])
        return plt.subplots(figsize=figsize, dpi=dpi)

    def _export_csv(self, df, metadata, stem):
        """
        Export chart data as CSV with metadata header rows.

        Args:
            df: The pandas DataFrame containing the chart data
            metadata: Chart metadata dictionary
            stem: Base filename for saving
        """
        csv_path = self.outdir / f"{stem}.csv"

        export_note = None
        if hasattr(df, "attrs"):
            export_note = df.attrs.get("export_note")

        # Create a copy of the data to avoid modifying the original
        csv_df = df.copy()

        # Prepare metadata rows
        meta_rows = []

        # Add author and generation info
        meta_rows.append(["Author", "Casey Dreier/The Planetary Society"])
        meta_rows.append(["Website", "https://planetary.org"])
        meta_rows.append(["Generated", datetime.now().strftime("%Y-%m-%d")])

        if "source" in metadata:
            meta_rows.append(["Data Source", metadata["source"]])
        if export_note:
            if isinstance(export_note, (list, tuple)):
                for note in export_note:
                    meta_rows.append(["Note", note])
            else:
                meta_rows.append(["Note", export_note])
        meta_rows.append(["License", "CC BY 4.0"])

        # Add a blank row between metadata and data
        meta_rows.append(["", ""])

        # Write metadata and data to CSV
        with open(csv_path, "w", newline="") as f:
            writer = csv.writer(f)

            # Write metadata rows
            for row in meta_rows:
                writer.writerow(row)

            # Write column names and data, converting NaN to empty strings,
            # formatting dates as YYYY-mm-dd, rounding floats to 2 significant
            # digits, and ensuring integers export as integers
            writer.writerow(csv_df.columns)
            for _, row in csv_df.iterrows():
                formatted_row = []
                for val in row:
                    if pd.isna(val):
                        formatted_row.append("")
                    elif hasattr(val, "strftime"):
                        # Format datetime/date/Timestamp as YYYY-mm-dd (no time)
                        formatted_row.append(val.strftime("%Y-%m-%d"))
                    elif isinstance(val, (int, float, np.integer, np.floating)):
                        if val == int(val):
                            formatted_row.append(int(val))
                        elif abs(val) > 1:
                            formatted_row.append(round(val, 2))
                        else:
                            formatted_row.append(val)
                    else:
                        formatted_row.append(val)
                writer.writerow(formatted_row)

        logger.debug(f"✓ saved {csv_path.name}")
        return csv_path

    def _get_fiscal_year_range_for_ticks(self, ax):
        """
        Extract fiscal year range from axis data for tick formatting.

        Default implementation for line charts that use continuous x-axis data.
        Subclasses can override this method to handle different data types
        (e.g., categorical data in bar charts).

        Args:
            ax: Matplotlib axes object

        Returns:
            tuple: (start_year, end_year, year_range) or None if not applicable
        """
        try:
            # Use dataLim (actual data extent) rather than get_xlim() which
            # includes matplotlib's auto-padding.  The padded range can push
            # the year_range into a sparser tick-density branch (e.g. 29 yr
            # data appearing as 32 yr range), producing inconsistent labels.
            dlim = ax.dataLim
            x0, x1 = dlim.x0, dlim.x1

            try:
                import matplotlib.dates as mdates

                start_year = mdates.num2date(x0).year
                end_year = mdates.num2date(x1).year
                year_range = abs(end_year - start_year)
                return (start_year, end_year, year_range)
            except Exception:
                # Fallback: assume values are years directly
                start_year = int(x0)
                end_year = int(x1)
                year_range = abs(end_year - start_year)
                return (start_year, end_year, year_range)

        except Exception:
            return None

    def _apply_fiscal_year_ticks(self, ax, style, tick_size=None):
        """
        Apply consistent fiscal year tick formatting to the x-axis.

        Sets major ticks at decade boundaries (years ending in 0),
        minor ticks at each year, and formats all labels horizontally.

        Args:
            ax: Matplotlib axes object
            style: dict of MOBILE or DESKTOP style options
            tick_size: Optional font size for tick labels
        """

        # Set major ticks at decade boundaries (years divisible by 10)
        ax.xaxis.set_major_locator(mdates.YearLocator(5))  # Every 5 years
        ax.xaxis.set_minor_locator(mdates.YearLocator(1))  # Every year

        # Format to show only the year
        ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y"))

        # Get year range using overrideable method
        year_info = self._get_fiscal_year_range_for_ticks(ax)

        if year_info is None:
            # Fallback to standard formatting if year range can't be determined
            if tick_size is None:
                tick_size = style.get("tick_size")
            plt.setp(
                ax.get_xticklabels(), rotation=style.get("tick_rotation", 0), fontsize=tick_size
            )
            return ax

        _start_year, _end_year, year_range = year_info

        # Adjust tick density based on year range
        decade_threshold = style.get("decade_tick_threshold", 50)
        if year_range > decade_threshold:
            # Very long ranges: decade labels only
            def decade_label(year, pos):
                year_int = int(mdates.num2date(year).year)
                return str(year_int) if year_int % 10 == 0 else ""

            ax.xaxis.set_major_formatter(FuncFormatter(decade_label))
        elif year_range > 20:
            # 20-30 years: show every 5 years
            ax.xaxis.set_major_locator(mdates.YearLocator(5))
        elif year_range > 10:
            # 10-20 years: show every 2 years for better readability
            ax.xaxis.set_major_locator(mdates.YearLocator(2))
        else:
            # Less than 10 years: show all years
            ax.xaxis.set_major_locator(mdates.YearLocator(1))

        # Minor ticks render at zero length/width by default (house style);
        # a style dict can re-enable them by setting the length/width keys.
        # Scoped to the x-axis: y tick marks are suppressed globally (the
        # hairline grid does that job) and must not be resurrected here.
        ax.tick_params(
            axis="x",
            which="minor",
            length=style.get("minor_tick_length", 4),
            width=style.get("minor_tick_width", 1),
        )
        ax.tick_params(
            axis="x",
            which="major",
            length=style.get("major_tick_length", 8),
            width=style.get("major_tick_width", 1.2),
        )

        # Allow override on tick size
        if tick_size is None:
            tick_size = style.get("tick_size")

        # Set tick labels horizontal and apply font size if provided
        plt.setp(ax.get_xticklabels(), rotation=style.get("tick_rotation", 0), fontsize=tick_size)

        return ax

    def _convert_xlim_to_datetime(self, xlim, x_data):
        """
        Convert xlim integer year values to datetime when x_data contains dates.

        When fiscal year data is converted to datetime objects, xlim values
        specified as integers (e.g., [2006, 2026]) must also be converted
        to datetime for matplotlib to properly clip the axis.

        Args:
            xlim: X-axis limits as tuple (min, max), dict, or None
            x_data: The x-axis data to check for datetime type

        Returns:
            Converted xlim (tuple of datetime or original value)
        """
        if xlim is None or x_data is None:
            return xlim

        # Only convert if x_data contains datetime objects
        if not self._contains_dates(x_data):
            return xlim

        # Check if x_data actually has datetime objects (not just integer years)
        try:
            first_elem = x_data.iloc[0] if hasattr(x_data, "iloc") else x_data[0]
            # Only convert if x_data has actual datetime objects
            if not (hasattr(first_elem, "year") and hasattr(first_elem, "month")):
                return xlim
        except (KeyError, IndexError):
            return xlim

        from datetime import datetime

        def convert_value(val):
            """Convert a single value to datetime if it's an integer year."""
            if isinstance(val, int) and 1900 <= val <= 2100:
                return datetime(val, 1, 1)
            if isinstance(val, float) and 1900 <= val <= 2100:
                return datetime(int(val), 1, 1)
            return val

        if isinstance(xlim, dict):
            # Handle dict format like {'left': 2006, 'right': 2026}
            return {k: convert_value(v) for k, v in xlim.items()}
        elif isinstance(xlim, (list, tuple)) and len(xlim) == 2:
            # Handle tuple/list format like [2006, 2026]
            return (convert_value(xlim[0]), convert_value(xlim[1]))

        return xlim

    # Helper to detect if x_data contains dates
    def _contains_dates(self, x_data):
        """
        Check if x_data contains date-like objects.

        Args:
            x_data: The x-axis data to check

        Returns:
            bool: True if the data appears to contain dates
        """
        if x_data is None or len(x_data) == 0:
            return False

        # Check if x_data contains datetime objects
        try:
            first_elem = x_data.iloc[0] if hasattr(x_data, "iloc") else x_data[0]
        except KeyError:
            logger.warning(f"Cannot read first element in array to check date objects: {x_data}")
            return False

        # A transition-quarter timeline is an ordered set of fiscal-period
        # labels, not a datetime axis. Checking only its first value (usually a
        # four-digit year) would incorrectly apply matplotlib date locators.
        if any("TQ" in str(value).upper() for value in x_data):
            return False

        # Check for datetime-like objects (Python datetime, pandas Timestamp)
        if hasattr(first_elem, "year") and hasattr(first_elem, "month"):
            return True

        # Check for numpy datetime64 scalar directly
        if isinstance(first_elem, np.datetime64):
            return True

        # Check array dtype for datetime64 (handles numpy arrays and pandas dtypes)
        if hasattr(x_data, "dtype"):
            try:
                if np.issubdtype(x_data.dtype, np.datetime64):
                    return True
            except TypeError:
                # pandas extension dtypes (e.g., Int64Dtype) are not always NumPy-compatible
                pass

            if pd.api.types.is_datetime64_any_dtype(x_data.dtype):
                return True

        # Check for integer years (1980, 1990, etc.)
        if isinstance(first_elem, int) and 1900 <= first_elem <= 2100:
            return True

        # Check for string years ("1980", "1990", etc.)
        return bool(
            isinstance(first_elem, str) and first_elem.isdigit() and 1900 <= int(first_elem) <= 2100
        )

    @staticmethod
    def _compute_scale_decimals(axis_limits: tuple, factor: float) -> tuple[int, float | None]:
        """Compute decimal places from axis range and scale factor."""
        try:
            range_value = abs(axis_limits[1] - axis_limits[0]) / factor
            return (1 if range_value < 10 else 0), range_value
        except (ZeroDivisionError, TypeError):
            return 0, None

    def _apply_scale_formatter(
        self,
        ax,
        scale: str = "billions",
        axis: str = "y",
        decimals: int | None = None,
        prefix: str | None = "$",
        tick_format: str | None = None,
    ):
        """
        Apply scale formatting to axis.

        Args:
            ax: The matplotlib Axes object to format
            scale: Scale to apply ('billions', 'millions', 'thousands', 'percentage')
            axis: Which axis to format ('x', 'y', or 'both')
            decimals: Number of decimal places to display
            prefix: Prefix to add before the number (e.g., '$')
        """
        scales = {
            "billions": {"factor": 1e9, "suffix": "B"},
            "millions": {"factor": 1e6, "suffix": "M"},
            "thousands": {"factor": 1e3, "suffix": "K"},
            "percentage": {"factor": 0.01, "suffix": "%", "prefix": ""},
        }

        if scale not in scales:
            warnings.warn(f"Scale '{scale}' not recognized. No formatter applied.", stacklevel=2)
            return

        scale_info = scales[scale]
        factor = scale_info["factor"]
        suffix = scale_info.get("suffix", "")
        prefix = scale_info.get("prefix", prefix)

        # Determine the number of decimals for smaller scales
        range_value = None
        if decimals is None:
            if axis in ("y", "both"):
                decimals, range_value = self._compute_scale_decimals(ax.get_ylim(), factor)
            elif axis == "x":
                decimals, range_value = self._compute_scale_decimals(ax.get_xlim(), factor)
            else:
                decimals = 0

        def formatter(x, pos):
            try:
                if not np.isfinite(x):
                    return ""
                if x == 0:
                    return ""
                if factor == 0:
                    return ""
                scaled_value = x / factor
                if tick_format:
                    formatted_num = format(scaled_value, tick_format)
                else:
                    # If range_value < 10, only show whole numbers
                    if (
                        axis in ("y", "both")
                        and decimals == 1
                        and range_value is not None
                        and range_value < 10
                    ):
                        if not np.isclose(scaled_value, round(scaled_value)):
                            return ""
                        format_spec = ".0f"
                    else:
                        format_spec = f".{decimals}f"
                    formatted_num = f"{scaled_value:{format_spec}}"
                return f"{prefix}{formatted_num}{suffix}"
            except Exception as e:
                logger.error(f"Formatter error for value x={x}, pos={pos}: {e}")
                return ""

        if axis in ("y", "both"):
            ax.yaxis.set_major_formatter(FuncFormatter(formatter))
        if axis in ("x", "both"):
            ax.xaxis.set_major_formatter(FuncFormatter(formatter))

    # One cap-height ascender + one descender. Every header line is measured
    # against this FIXED reference instead of its own glyphs, so the header's
    # vertical rhythm depends only on line count, font size, line spacing and dpi
    # — never on which characters a title/subtitle happens to contain. Without
    # this, a subtitle full of gerund descenders ("comparing spending") and an
    # all-caps one ("NASA'S BUDGET") could reserve different heights, giving the
    # inconsistent line-heights reported by users.
    _METRIC_REFERENCE_GLYPHS: ClassVar[str] = "Ag"

    def _measure_text_height(self, fig, text, fontsize, *, linespacing=1.2) -> float:
        """
        Deterministic, glyph-independent text-block height in figure fractions.

        The height is measured from a fixed reference string (:pydata:`_METRIC_REFERENCE_GLYPHS`
        per line) that mirrors ``text``'s line count — not from ``text`` itself —
        so descenders, all-caps, brackets or accents can never shift the header
        geometry. Only line count, ``fontsize``, ``linespacing`` and dpi matter.

        Args:
            fig: The matplotlib Figure object
            text: The text whose line count drives the measurement
            fontsize: Font size for the text
            linespacing: Spacing between lines as a multiple of the font size

        Returns:
            float: Height of the text block in figure-fraction coordinates (0.0-1.0)
        """
        if not text:
            return 0.0

        # Map each real line to the fixed reference (blank lines stay blank so
        # matplotlib's zero-height handling for empty lines is preserved). This
        # makes the measured block height a pure function of the line pattern.
        lines = str(text).splitlines() or [str(text)]
        reference = "\n".join(
            self._METRIC_REFERENCE_GLYPHS if line.strip() else "" for line in lines
        )

        # The height is a pure function of these inputs — memoize so repeated
        # header measurement/render passes skip the throwaway probe artist.
        cache_key = (reference, fontsize, linespacing, fig.get_figheight(), fig.dpi)
        if not hasattr(self, "_text_height_cache"):
            self._text_height_cache: dict = {}
        cached = self._text_height_cache.get(cache_key)
        if cached is not None:
            return cached

        # Create temporary invisible text element to measure
        temp_text = fig.text(
            0,
            0,
            reference,
            fontsize=fontsize,
            linespacing=linespacing,
            alpha=0.0,
        )

        renderer = fig.canvas.get_renderer()

        # Get bounding box in display (pixel) coordinates
        bbox = temp_text.get_window_extent(renderer)

        # Convert to figure-fraction coordinates
        fig_height_px = fig.get_figheight() * fig.dpi
        height = bbox.height / fig_height_px

        # Clean up the temporary text
        temp_text.remove()

        self._text_height_cache[cache_key] = height
        return height

    def _wrap_header_text(
        self,
        fig,
        text: str | None,
        fontsize: float,
        *,
        fontweight: str | int | None = None,
        left_margin: float = 0.01,
        right_margin: float = 0.99,
        fallback_wrap_length: int = 65,
    ) -> str:
        """
        Wrap header/subtitle text to fit rendered pixel width between figure margins.

        This avoids character-count wrapping artifacts where proportional fonts
        can overflow the right edge even when the wrapped line length looks safe.

        ``fontweight`` should match how the text will actually be drawn (e.g.
        ``"bold"`` for the title); bold glyphs are wider than the default weight,
        so measuring with the render weight keeps wrapped lines inside the margins.
        """
        if text is None:
            return ""

        escaped_text = self._escape_svg_text(text)
        if not escaped_text:
            return ""

        max_width_px = (right_margin - left_margin) * fig.get_figwidth() * fig.dpi
        if max_width_px <= 0:
            return "\n".join(
                textwrap.wrap(str(escaped_text), width=fallback_wrap_length) or [str(escaped_text)]
            )

        # Wrapping is deterministic in (text, size, weight, pixel width), and both
        # the header-measurement and header-render passes wrap the same strings —
        # memoize so the second pass skips the full fig.canvas.draw() below.
        cache_key = (str(escaped_text), fontsize, fontweight, round(max_width_px, 3))
        if not hasattr(self, "_wrap_text_cache"):
            self._wrap_text_cache: dict = {}
        cached = self._wrap_text_cache.get(cache_key)
        if cached is not None:
            return cached

        fig.canvas.draw()
        renderer = fig.canvas.get_renderer()
        wrapped_lines: list[str] = []
        probe_kwargs = {"fontsize": fontsize, "alpha": 0.0}
        if fontweight is not None:
            probe_kwargs["fontweight"] = fontweight
        probe = fig.text(0, 0, "", **probe_kwargs)

        try:
            paragraphs = str(escaped_text).splitlines() or [""]
            for paragraph in paragraphs:
                words = paragraph.split()
                if not words:
                    wrapped_lines.append("")
                    continue

                current_line = words[0]
                for word in words[1:]:
                    candidate = f"{current_line} {word}"
                    probe.set_text(candidate)
                    if probe.get_window_extent(renderer).width <= max_width_px:
                        current_line = candidate
                    else:
                        wrapped_lines.append(current_line)
                        current_line = word

                wrapped_lines.append(current_line)
        finally:
            probe.remove()

        wrapped = "\n".join(wrapped_lines)
        self._wrap_text_cache[cache_key] = wrapped
        return wrapped

    def _eyebrow_text(self, metadata, style) -> str | None:
        """Return the uppercased/escaped eyebrow string, or None when it won't render.

        The eyebrow is desktop-only by doctrine: it renders only when the device
        style enables ``show_eyebrow`` *and* an ``eyebrow`` value is present.
        """
        if not style.get("show_eyebrow"):
            return None
        eyebrow = metadata.get("eyebrow")
        if not eyebrow:
            return None
        return self._escape_svg_text(str(eyebrow).upper())

    @staticmethod
    def _eyebrow_size(style) -> float:
        """Eyebrow font size — one rule shared by measurement and rendering."""
        return style["title_size"] * tokens.EYEBROW_SIZE_RATIO

    def _wrap_title(self, fig, title: str | None, style) -> str:
        """Wrap the bold title to the header margins (measured at bold weight).

        Mirrors the subtitle wrapping so long titles (notably on the narrow
        mobile 8:9 canvas) break onto multiple lines instead of running off the
        right edge. For titles that already fit on one line this returns the
        title unchanged, keeping existing output byte-identical.
        """
        header_x = style.get("header_x", 0.01)
        return self._wrap_header_text(
            fig,
            title,
            style["title_size"],
            fontweight="bold",
            left_margin=header_x,
            right_margin=1 - header_x,
            fallback_wrap_length=style.get("subtitle_wrap_length", 65),
        )

    def _calculate_header_height(self, fig, metadata, style) -> float:
        """
        Calculate header height based on actual content dimensions.

        Measures the (wrapped) title, wrapped subtitle, and optional eyebrow to
        determine the exact header space needed, avoiding both wasted space and
        text overlap.

        Returns:
            header_height in figure-fraction coordinates.
        """
        title = metadata.get("title")
        subtitle = metadata.get("subtitle")
        eyebrow = self._eyebrow_text(metadata, style)

        if not title and not subtitle and not eyebrow:
            return 0.0

        # Measure eyebrow height (single line above the title). Includes the
        # kicker gap so the reserved header height matches what _add_header
        # draws — otherwise the eyebrow's stack eats the clearance between the
        # subtitle and the plot.
        eyebrow_height = 0.0
        if eyebrow:
            eyebrow_size = self._eyebrow_size(style)
            eyebrow_height = (
                self._measure_text_height(fig, eyebrow, eyebrow_size) * tokens.EYEBROW_STACK_RATIO
            )

        # Measure title height (wrapped, so multi-line titles get header room)
        title_line_spacing = style.get("title_line_spacing", 1.1)
        wrapped_title = self._wrap_title(fig, title, style)
        title_height = self._measure_text_height(
            fig, wrapped_title, style["title_size"], linespacing=title_line_spacing
        )

        # Measure subtitle height (with wrapping applied)
        subtitle_ratio = style.get("subtitle_size_ratio", 0.7)
        subtitle_size = style["title_size"] * subtitle_ratio
        subtitle_line_spacing = style.get("subtitle_line_spacing", 1.05)
        header_x = style.get("header_x", 0.01)
        if subtitle:
            wrapped = self._wrap_header_text(
                fig,
                subtitle,
                subtitle_size,
                left_margin=header_x,
                right_margin=1 - header_x,
                fallback_wrap_length=style.get("subtitle_wrap_length", 65),
            )
            subtitle_height = self._measure_text_height(
                fig,
                wrapped,
                subtitle_size,
                linespacing=subtitle_line_spacing,
            )
        else:
            subtitle_height = 0.0

        # Add padding: top margin + gap between title/subtitle + bottom margin
        padding = style.get("header_padding", 0.03)

        # Get min/max constraints from style (with sensible defaults)
        min_height = style.get("header_min_height", 0.06)
        max_height = style.get("header_max_height", 0.18)

        # Calculate total height and constrain to bounds
        total_height = eyebrow_height + title_height + subtitle_height + padding
        return max(min_height, min(max_height, total_height))

    def _add_header(self, fig, metadata, style):
        """
        Add header elements to the figure: title and subtitle with left alignment.

        Subtitle is positioned dynamically relative to the title's actual height,
        rather than using a hardcoded y-offset.

        Args:
            fig: The matplotlib Figure object
            metadata: Chart metadata dictionary
            style: Style dictionary (DESKTOP or MOBILE)
        """
        title = metadata.get("title")
        subtitle = metadata.get("subtitle")
        eyebrow = self._eyebrow_text(metadata, style)

        header_x = style.get("header_x", 0.01)
        header_y = style.get("header_y", 0.98)
        subtitle_ratio = style.get("subtitle_size_ratio", 0.7)
        subtitle_line_spacing = style.get("subtitle_line_spacing", 1.05)
        gap = style.get("title_subtitle_gap", 0.005)

        # Top of the title zone. When an eyebrow is present it renders at the very
        # top and pushes the title (and thus the subtitle) down beneath it.
        title_top_y = header_y
        if eyebrow:
            eyebrow_size = self._eyebrow_size(style)
            # Optional signature rule just above the eyebrow, in the top margin
            # (above header_y), so it needs no reserved header space.
            if tokens.EYEBROW_RULE_COLOR is not None:
                fig.add_artist(
                    plt.Line2D(
                        [header_x, header_x + tokens.EYEBROW_RULE_LENGTH],
                        [header_y + 0.006, header_y + 0.006],
                        transform=fig.transFigure,
                        color=tokens.EYEBROW_RULE_COLOR,
                        linewidth=tokens.EYEBROW_RULE_LINEWIDTH,
                        solid_capstyle="butt",
                    )
                )
            fig.text(
                header_x,
                header_y,
                eyebrow,
                fontsize=eyebrow_size,
                fontweight=tokens.EYEBROW_WEIGHT,
                color=tokens.EYEBROW_COLOR,
                ha="left",
                va="top",
            )
            # Deterministic: the eyebrow is top-anchored at header_y, so the
            # title sits a glyph-independent block height below it. Only a
            # fraction of the measured "Ag" box is reserved (EYEBROW_STACK_RATIO)
            # so the title climbs into the kicker's unused descender whitespace
            # and the two read as one headline block.
            eyebrow_height = self._measure_text_height(fig, eyebrow, eyebrow_size)
            title_top_y = header_y - eyebrow_height * tokens.EYEBROW_STACK_RATIO

        # Track vertical position (starts at top of the title zone)
        title_bottom_y = title_top_y  # Default if no title
        title_text = subtitle_text = None

        # Add title if provided (wrapped to the header margins). The tight
        # title line spacing keeps a wrapped multi-line title's line-height in
        # rhythm with the rest of the header instead of matplotlib's loose 1.2.
        if title:
            title_line_spacing = style.get("title_line_spacing", 1.1)
            wrapped_title = self._wrap_title(fig, title, style)
            title_text = fig.text(
                header_x,
                title_top_y,
                wrapped_title,
                fontsize=style["title_size"],
                fontweight="bold",
                color=tokens.TITLE_COLOR,
                ha="left",
                va="top",
                linespacing=title_line_spacing,
            )

            # Deterministic subtitle anchor: the title is top-anchored at
            # title_top_y, so its bottom edge is that anchor minus the
            # glyph-independent block height (no dependence on the title's ink).
            title_height = self._measure_text_height(
                fig, wrapped_title, style["title_size"], linespacing=title_line_spacing
            )
            title_bottom_y = title_top_y - title_height - gap

        # Add subtitle if provided, positioned relative to title
        if subtitle:
            subtitle_size = style["title_size"] * subtitle_ratio
            wrapped_subtitle = self._wrap_header_text(
                fig,
                subtitle,
                subtitle_size,
                left_margin=header_x,
                right_margin=1 - header_x,
                fallback_wrap_length=style.get("subtitle_wrap_length", 65),
            )
            subtitle_text = fig.text(
                header_x,
                title_bottom_y,
                wrapped_subtitle,
                fontsize=subtitle_size,
                linespacing=subtitle_line_spacing,
                color=tokens.SUBTITLE_COLOR,
                ha="left",
                va="top",
            )

        return title_text, subtitle_text

    def _center_subtitle_vertically(self, fig, title_text, subtitle_text, style, header_height):
        """Nestle the subtitle under the title, growing the plot up to fill header slack.

        Places the subtitle with scaled padding under the title, then grows the
        plot's TOP edge up into the leftover header space while keeping its
        bottom anchored — so the plot stays filled toward the footer instead of
        vacating a band of dead space there (the old translate-up behaviour).
        """
        if title_text is None or subtitle_text is None:
            return

        visible_axes = self._visible_axes(fig)
        if not visible_axes:
            return

        # Deterministic, glyph-independent geometry. The title is top-anchored,
        # so its bottom edge is its anchor y minus a block height that depends
        # only on line count / size / spacing; the subtitle block height is
        # likewise glyph-independent. This keeps the subtitle's vertical rhythm
        # identical regardless of descenders or all-caps content.
        title_top_y = title_text.get_position()[1]
        title_bottom_y = title_top_y - self._measure_text_height(
            fig,
            title_text.get_text(),
            title_text.get_fontsize(),
            linespacing=title_text.get_linespacing(),
        )
        subtitle_height = self._measure_text_height(
            fig,
            subtitle_text.get_text(),
            subtitle_text.get_fontsize(),
            linespacing=subtitle_text.get_linespacing(),
        )
        chart_top_y = max(ax.get_position().y1 for ax in visible_axes)
        available_spacing = title_bottom_y - chart_top_y - subtitle_height
        if available_spacing <= 0:
            return

        padding_scale = style.get("subtitle_vertical_padding_scale", 1.0)
        target_spacing = available_spacing * padding_scale
        requested_shift = available_spacing - target_spacing
        header_boundary_y = 1.0 - header_height
        available_upward_shift = max(0.0, header_boundary_y - chart_top_y)
        axes_shift = min(requested_shift, available_upward_shift)

        # Grow the plot's TOP edge up into the header slack, keeping its bottom
        # anchored — the plot already fills the footer side after
        # _stretch_axes_vertically, so translating the whole plot up (the old
        # behaviour) only vacated an equal band of dead space above the footer.
        # Bottom-anchored growth reclaims the slack while keeping the plot filled.
        if axes_shift > 0:
            axes_bottom = min(ax.get_position().y0 for ax in visible_axes)
            axes_top = max(ax.get_position().y1 for ax in visible_axes)
            axes_height = axes_top - axes_bottom
            if axes_height > 0:
                self._remap_axes_band(
                    visible_axes,
                    "y",
                    axes_bottom,
                    axes_height,
                    axes_bottom,
                    axes_height + axes_shift,
                )

        # Deterministic title→subtitle gap (matches the eyebrow→title rhythm).
        # `title_bottom_y` is the title's glyph-independent box bottom — the
        # DESCENDER line, well below the baseline where the title's body of text
        # visually sits. Anchoring the subtitle there makes a title whose last
        # line has a descender ("...Fundin[g]") read as a big gap. So pull the
        # subtitle up by a fraction of the TITLE line box, reclaiming the
        # title's descender depth so the subtitle hugs the baseline instead.
        # Basing the reclaim on the title (not subtitle) box is what keeps this
        # consistent: the descender being reclaimed is the title's own.
        title_line_height = self._measure_text_height(
            fig,
            self._METRIC_REFERENCE_GLYPHS,
            title_text.get_fontsize(),
            linespacing=title_text.get_linespacing(),
        )
        overlap = title_line_height * tokens.TITLE_SUBTITLE_OVERLAP_RATIO
        subtitle_text.set_y(title_bottom_y + overlap)

    def _add_footer(self, fig, metadata, style, bottom_margin):
        """
        Add footer elements to the figure: horizontal line, source text, and logo.

        Args:
            fig: The matplotlib Figure object
            metadata: Chart metadata dictionary
            style: Style dictionary (DESKTOP or MOBILE)
            bottom_margin: Bottom margin to reserve for the footer
        """
        # NOTE: Caller (_adjust_layout_for_header_footer) already checks whether
        # footer should be displayed based on style and metadata settings.
        # Layout spacing is handled by tight_layout(rect=...) in _adjust_layout_for_header_footer,
        # so we do NOT call fig.subplots_adjust here to avoid conflicting layout calls.

        # Optional full-width footer band behind the footer zone (theme seam;
        # None = no band, the default).
        if tokens.FOOTER_BAND_COLOR is not None:
            fig.add_artist(
                plt.Rectangle(
                    (0, 0),
                    1,
                    bottom_margin,
                    transform=fig.transFigure,
                    facecolor=tokens.FOOTER_BAND_COLOR,
                    edgecolor="none",
                    zorder=1,
                )
            )

        # Add horizontal spacer line, matching header text margins
        spacer_y = bottom_margin / 2  # Place line halfway in the margin
        footer_color = style.get("footer_color", tokens.FOOTER_COLOR)
        self._add_horizontal_spacer(
            fig,
            y_position=spacer_y,
            color=footer_color,
            linewidth=style.get("footer_line_width", 1),
            extent=style.get("footer_extent", (0.01, 0.99)),
        )

        # Add source if provided
        source_text = metadata.get("source")
        if source_text:
            self._add_source(fig, source_text, style)

        # Add methodology note (if provided) directly above the source line
        note_text = metadata.get("note")
        if note_text:
            self._add_note(fig, note_text, style)

        # Add logo if enabled in the style
        if style.get("add_logo", True):
            self._add_logo(fig, style)

    def _adjust_layout_for_header_footer(self, fig, metadata, style):
        """
        Adjust figure layout to accommodate headers and footers.

        This is the SINGLE SOURCE OF TRUTH for layout spacing. Uses tight_layout(rect=...)
        to allocate space for header/footer, with dynamic header sizing based on actual
        content measurements.

        Args:
            fig: The matplotlib Figure object
            metadata: Chart metadata dictionary
            style: Style dictionary (DESKTOP or MOBILE, etc)
        """
        # Determine if header/footer should be displayed
        show_header = style.get("header") or metadata.get("header")
        show_footer = style.get("footer") or metadata.get("footer")

        # Calculate dynamic header height based on actual content
        if show_header:
            # Use dynamic calculation, falling back to style default
            header_height = self._calculate_header_height(fig, metadata, style)
            # Ensure we don't go below the style minimum
            style_min = style.get("header_height", 0.1)
            header_height = max(header_height, style_min)
        else:
            header_height = 0

        # Footer height is fixed (source text + logo don't vary much)
        footer_height = style.get("footer_height", 0) if show_footer else 0

        # Add header if enabled (after calculating height so it knows space available)
        title_text = subtitle_text = None
        if show_header:
            title_text, subtitle_text = self._add_header(fig, metadata, style)

        # Add footer if enabled
        if show_footer:
            self._add_footer(fig, metadata, style, footer_height)

        # Apply tight layout with adjusted rectangle
        # rect = [left, bottom, right, top] in figure-fraction coordinates
        fig.tight_layout(rect=[0, footer_height, 1, 1 - header_height])

        if show_header or show_footer:
            self._stretch_axes_vertically(fig, header_height, footer_height, style)
            self._align_axes_horizontally(fig, style)
        else:
            self._center_axes_vertically(fig, header_height, footer_height)
        self._center_subtitle_vertically(fig, title_text, subtitle_text, style, header_height)

        return fig

    def _visible_axes(self, fig):
        """Return the figure's visible axes."""
        return [ax for ax in fig.get_axes() if ax.get_visible()]

    def _visible_axes_fig_bboxes(self, fig, visible_axes, renderer):
        """Tight bboxes of the given axes (incl. labels) in figure-fraction coords."""
        fig_bboxes = []
        for ax in visible_axes:
            bbox = ax.get_tightbbox(renderer)
            if bbox is not None:
                fig_bboxes.append(bbox.transformed(fig.transFigure.inverted()))
        return fig_bboxes

    def _stretch_axes_vertically(self, fig, header_height, footer_height, style):
        """Expand rendered chart bounds to fill the usable header/footer zone."""
        visible_axes = self._visible_axes(fig)
        if not visible_axes:
            return

        padding = style.get("chart_vertical_padding", 0.01)
        target_bottom = footer_height + padding
        target_top = 1.0 - header_height - padding
        if target_top <= target_bottom:
            return

        for _ in range(2):
            fig.canvas.draw()
            renderer = fig.canvas.get_renderer()
            visual_bounds = self._visible_axes_fig_bboxes(fig, visible_axes, renderer)
            if not visual_bounds:
                return

            visual_bottom = min(bounds.y0 for bounds in visual_bounds)
            visual_top = max(bounds.y1 for bounds in visual_bounds)
            axes_bottom = min(ax.get_position().y0 for ax in visible_axes)
            axes_top = max(ax.get_position().y1 for ax in visible_axes)
            axes_height = axes_top - axes_bottom
            if axes_height <= 0:
                return

            aligned_bottom = axes_bottom + target_bottom - visual_bottom
            aligned_top = axes_top + target_top - visual_top
            aligned_height = aligned_top - aligned_bottom
            if aligned_height <= 0:
                return

            self._remap_axes_band(
                visible_axes, "y", axes_bottom, axes_height, aligned_bottom, aligned_height
            )

    @staticmethod
    def _remap_axes_band(visible_axes, axis, old_start, old_size, new_start, new_size):
        """Proportionally remap all axes from one figure-fraction band to another.

        Each axes' start/end within the old band keeps its relative position in
        the new band. ``axis`` is ``"y"`` (vertical remap) or ``"x"``.
        """
        for ax in visible_axes:
            pos = ax.get_position()
            lo, hi = (pos.y0, pos.y1) if axis == "y" else (pos.x0, pos.x1)
            new_lo = new_start + (lo - old_start) / old_size * new_size
            new_hi = new_start + (hi - old_start) / old_size * new_size
            if axis == "y":
                ax.set_position([pos.x0, new_lo, pos.width, new_hi - new_lo])
            else:
                ax.set_position([new_lo, pos.y0, new_hi - new_lo, pos.height])

    def _align_axes_horizontally(self, fig, style):
        """Align rendered chart edges with the header and footer text extent."""
        visible_axes = self._visible_axes(fig)
        if not visible_axes:
            return

        target_left, target_right = style.get("footer_extent", (0.01, 0.99))

        for _ in range(2):
            fig.canvas.draw()
            renderer = fig.canvas.get_renderer()
            visual_bounds = self._visible_axes_fig_bboxes(fig, visible_axes, renderer)
            if not visual_bounds:
                return

            visual_left = min(bounds.x0 for bounds in visual_bounds)
            visual_right = max(bounds.x1 for bounds in visual_bounds)
            axes_left = min(ax.get_position().x0 for ax in visible_axes)
            axes_right = max(ax.get_position().x1 for ax in visible_axes)
            axes_width = axes_right - axes_left
            if axes_width <= 0:
                return

            aligned_left = axes_left + target_left - visual_left
            aligned_right = axes_right + target_right - visual_right
            aligned_width = aligned_right - aligned_left

            self._remap_axes_band(
                visible_axes, "x", axes_left, axes_width, aligned_left, aligned_width
            )

    def _center_axes_vertically(self, fig, header_height, footer_height):
        """Shift axes to equalize spacing between header content and footer content."""
        visible_axes = self._visible_axes(fig)
        if not visible_axes:
            return

        renderer = fig.canvas.get_renderer()

        # Measure visual extent using get_tightbbox (includes tick labels, axis labels)
        fig_bboxes = self._visible_axes_fig_bboxes(fig, visible_axes, renderer)
        if not fig_bboxes:
            return

        visual_top = max(bbox.y1 for bbox in fig_bboxes)
        visual_bottom = min(bbox.y0 for bbox in fig_bboxes)

        # Reference boundaries for visual centering:
        # - Top: header zone bottom (subtitle bottom edge)
        # - Bottom: footer spacer line position (visual top of footer content)
        ref_top = 1.0 - header_height
        ref_bottom = footer_height / 2.0 if footer_height > 0 else 0.0

        top_gap = ref_top - visual_top
        bottom_gap = visual_bottom - ref_bottom

        # Positive shift = move chart down (subtract from y-coords)
        shift = (bottom_gap - top_gap) / 2.0

        if abs(shift) < 0.005:
            return

        # Safety: don't push visual extent outside the tight_layout zone
        zone_top = 1.0 - header_height
        zone_bottom = footer_height
        new_visual_top = visual_top - shift
        new_visual_bottom = visual_bottom - shift

        if new_visual_top > zone_top or new_visual_bottom < zone_bottom:
            if shift > 0:  # moving down
                shift = min(shift, visual_bottom - zone_bottom)
            else:  # moving up
                shift = max(shift, -(zone_top - visual_top))
            if abs(shift) < 0.005:
                return
            logger.debug("Chart centering shift clamped to %.4f", shift)

        for ax in visible_axes:
            pos = ax.get_position()
            ax.set_position([pos.x0, pos.y0 - shift, pos.width, pos.height])

    def _add_horizontal_spacer(
        self, fig, y_position=None, color=None, linewidth=0.5, extent=(0.02, 0.98)
    ):
        """
        Add a horizontal line spacer to the figure.

        Args:
            fig: The matplotlib Figure object
            y_position: Y-position of the line in figure coordinates
            color: Color of the line
            linewidth: Width of the line
            extent: Tuple of (start, end) x-positions in figure coordinates
        """
        # Set default values if not provided
        if y_position is None:
            y_position = 0.06

        if color is None:
            color = tokens.FOOTER_COLOR

        # Add the horizontal line
        fig.add_artist(
            plt.Line2D(
                [extent[0], extent[1]],  # x-positions (left, right)
                [y_position, y_position],  # y-positions (same for horizontal line)
                transform=fig.transFigure,  # Use figure coordinates
                color=color,
                linestyle="-",
                linewidth=linewidth,
            )
        )

    def _add_logo(self, fig, style):
        """Add The Planetary Society vector logo to the figure footer."""
        from tpsplots.views.logo import draw_logo

        try:
            draw_logo(
                fig,
                zoom=style.get("logo_zoom", 0.03),
                x=style.get("logo_x", 0.01),
                y=style.get("logo_y", 0.01),
                color=tokens.LOGO_COLOR,
            )
        except (FileNotFoundError, OSError, ValueError, struct.error) as e:
            logger.error(f"Warning: Could not add logo: {e}")

    def _add_source(self, fig, source_text, style):
        """
        Add source text to the bottom right of the figure.

        Args:
            fig: The matplotlib Figure object
            source_text: Source text to display
            style: Style dictionary (DESKTOP or MOBILE)
        """
        if not source_text:
            return

        # Add text at the bottom right, matching footer line right edge
        footer_extent = style.get("footer_extent", (0.01, 0.99))
        fig.text(
            footer_extent[1],  # x position (right side, matches footer line)
            style.get("source_y", 0.01),  # y position (bottom)
            f"Source: {source_text}".upper(),
            fontsize=style.get("source_size", 11),
            color=style.get("footer_color", tokens.FOOTER_COLOR),
            ha="right",
            va="bottom",
        )

    def _add_note(self, fig, note_text, style):
        """Add a right-aligned italic methodology note above the source line.

        Single line, no wrapping (by design). Positioned one source-line-height
        above the source baseline so it clears the source attribution text.

        Args:
            fig: The matplotlib Figure object
            note_text: Methodology note to display
            style: Style dictionary (DESKTOP or MOBILE)
        """
        if not note_text:
            return

        source_size = style.get("source_size", 11)
        note_size = source_size * 0.9
        source_y = style.get("source_y", 0.01)
        footer_extent = style.get("footer_extent", (0.01, 0.99))

        # Reserve one source-line of vertical space so the note sits above it.
        line_height = self._measure_text_height(fig, "Ag", source_size)

        fig.text(
            footer_extent[1],  # right edge, matching the source line
            source_y + line_height * 1.2,
            self._escape_svg_text(note_text),
            fontsize=note_size,
            color=style.get("footer_color", tokens.FOOTER_COLOR),
            style="italic",
            ha="right",
            va="bottom",
        )

    def _save_chart(self, fig, filename, metadata, create_pptx=False, create_svg=True) -> list[str]:
        """
        Save chart as PNG, and optionally SVG and PPTX.

        Args:
            fig: The matplotlib Figure object
            filename: Base filename for saving
            metadata: title, source, etc context for chart
            create_pptx: Whether to create a PowerPoint file
            create_svg: Whether to create an SVG file (default True)

        Returns:
            list[str]: File paths created by this save operation
        """

        clean_filename = (
            filename.replace("_desktop", "").replace("_mobile", "").replace("_social", "")
        )
        png_path = self.outdir / f"{filename}.png"

        file_metadata = {
            "Title": metadata.get("title", clean_filename.replace("_", " ")),
            "Creator": "Casey Dreier/The Planetary Society",
            "Date": datetime.today().strftime("%Y-%m-%d"),
            "Rights": "CC BY 4.0",
            "Source": metadata.get("source"),
        }

        files: list[str] = []

        if create_svg:
            svg_path = self.outdir / f"{filename}.svg"
            fig.savefig(svg_path, metadata=file_metadata, format="svg", dpi="figure")
            files.append(str(svg_path))

        fig.savefig(png_path, metadata=file_metadata, format="png", dpi="figure")
        files.append(str(png_path))
        logger.debug(f"✓ saved {', '.join(Path(p).name for p in files)}")

        if create_pptx:
            pptx_path = self.outdir / f"{filename.replace('_desktop', '')}.pptx"
            self._create_pptx(png_path, pptx_path, metadata)
            logger.debug(f"✓ saved {pptx_path.name}")
            files.append(str(pptx_path))

        plt.close(fig)
        return files

    def _create_pptx(self, png_path, pptx_path, metadata=None):
        """
        Create a PowerPoint file with the chart, scaled by height to fit completely in a 16x9 slide.

        Args:
            png_path: Path to the PNG image to include
            pptx_path: Path for the output PowerPoint file
        """
        from PIL import Image
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches

        if metadata is None:
            metadata = {}
        prs = Presentation()
        # Set slide size to 16x9 (in inches)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        # Set slide background color to self.TPS_COLORS['Slushy Brine']

        bg_color = self.TPS_COLORS["Slushy Brine"]
        # Convert hex bg_color to RGB
        if isinstance(bg_color, str) and bg_color.startswith("#"):
            bg_color = bg_color.lstrip("#")
            bg_color = RGBColor(
                int(bg_color[0:2], 16), int(bg_color[2:4], 16), int(bg_color[4:6], 16)
            )

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Get image size in inches
        with Image.open(png_path) as img:
            img_width_px, img_height_px = img.size
            img_dpi = img.info.get("dpi", (300, 300))[0]
            img_width_in = img_width_px / img_dpi
            img_height_in = img_height_px / img_dpi

        # Scale image by height to fit slide
        target_height_in = prs.slide_height / Inches(1)
        scale = target_height_in / img_height_in
        scaled_width_in = img_width_in * scale
        scaled_height_in = img_height_in * scale

        # Center the image horizontally
        left = (prs.slide_width - Inches(scaled_width_in)) / 2
        top = 0  # Top align

        slide.shapes.add_picture(
            str(png_path), left, top, width=Inches(scaled_width_in), height=Inches(scaled_height_in)
        )

        # Prepare title, subtitle, and source text
        notes = []
        notes.append(metadata.get("title"))
        notes.append(metadata.get("subtitle"))
        notes.append("\n")
        if "source" in metadata:
            notes.append("Source: " + metadata.get("source", ""))
        notes.append(
            f"Author: {metadata.get('Creator', 'Casey Dreier/The Planetary Society')}\nGenerated: {datetime.now().strftime('%Y-%m-%d')}"
        )
        notes.append("License: CC BY 4.0")

        # Clear notes of None or empty strings
        notes = [note for note in notes if note and note.strip()]
        if notes:
            # Add to text frame
            # Add title and source text to the notes section of the slide
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = "\n".join(notes)

        # Save
        prs.save(pptx_path)

    def _parse_newlines_in_labels(self, labels):
        """
        Parse literal newline sequences in category labels to actual newlines.

        This allows users to include multiline labels in their data sources (CSV files,
        YAML configurations, etc.) by using the literal string "\\n". This method converts
        those literal sequences to actual newline characters that matplotlib will render
        as line breaks in tick labels.

        Args:
            labels: Array-like of label strings

        Returns:
            numpy array: Labels with parsed newlines

        Examples:
            >>> labels = ["Program\\nName", "Mission\\nType", "Single"]
            >>> parsed = self._parse_newlines_in_labels(labels)
            >>> parsed[0]
            'Program\nName'  # Actual newline character, not literal \\n
        """
        if labels is None:
            return labels

        # Convert to list if needed for processing
        labels_list = list(labels) if not isinstance(labels, list) else labels

        # Replace literal \n with actual newline character
        parsed_labels = [str(label).replace("\\n", "\n") for label in labels_list]

        return np.array(parsed_labels)

    def _escape_svg_text(self, text):
        """
        Escape special characters for SVG text rendering in matplotlib.

        Args:
            text: The text string to escape

        Returns:
            The escaped text string
        """
        if text is None:
            return None

        # Define replacements for special characters
        replacements = {"$": r"\$"}

        # Apply all replacements
        for char, replacement in replacements.items():
            text = text.replace(char, replacement)

        return text

    def _format_value(self, value, format_type: str) -> str:
        """
        Format values according to the specified format type.

        Args:
            value: The numeric value to format
            format_type: Format preset ('monetary', 'percentage', 'integer', 'float')
                        or a Python format spec (e.g., '.1f', '.2f', ',.0f')

        Returns:
            Formatted string representation of the value

        Raises:
            ValueError: If format_type is not recognized
        """
        import pandas as pd

        if pd.isna(value):
            return ""

        if format_type == "monetary":
            return self._format_monetary(value)
        elif format_type == "percentage":
            return f"{value:.1f}%"
        elif format_type == "integer":
            return f"{int(value):,}"
        elif format_type == "float":
            return f"{value:.1f}"
        else:
            # Try as custom Python format specification
            try:
                return f"{value:{format_type}}"
            except (ValueError, KeyError) as e:
                raise ValueError(
                    f"Invalid value_format: '{format_type}'. "
                    f"Must be one of 'monetary', 'percentage', 'integer', 'float' "
                    f"or a valid Python format spec (e.g., '.1f', '.2f', ',.0f'). "
                    f"Error formatting value {value}: {e}"
                ) from e

    def _format_monetary(self, value) -> str:
        """
        Format monetary values with appropriate suffixes (B, M, K).

        Args:
            value: The numeric value to format

        Returns:
            Formatted string with currency symbol and magnitude suffix
        """
        abs_value = abs(value)
        sign = "-" if value < 0 else ""

        if abs_value >= 1_000_000_000:
            return f"{sign}${abs_value / 1_000_000_000:.1f}B"
        elif abs_value >= 1_000_000:
            return f"{sign}${abs_value / 1_000_000:.0f}M"
        elif abs_value >= 1_000:
            return f"{sign}${abs_value / 1_000:.0f}K"
        else:
            return f"{sign}${abs_value:.0f}"
