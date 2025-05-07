from matplotlib.offsetbox import OffsetImage, AnnotationBbox
import matplotlib.image as mpimg
import numpy as np
from pathlib import Path
import matplotlib.pyplot as plt
from pywaffle import Waffle


class ChartView:
    """View component for chart generation with desktop/mobile versions built in."""
    
    # Load global style once at class initialization
    plt.style.use(Path(__file__).parent / "style" / "tps_base.mplstyle")
    
    # Shared color palette
    COLORS = {
        "blue": "#037CC2",
        "purple": "#643788",
        "orange": "#FF5D47",
        "light_blue": "#80BDE0",
        "light_purple": "#B19BC3",
        "lunar_dust": "#8C8C8C",
        "dark_gray": "#414141",
        "medium_gray": "#C3C3C3",
        "light_gray": "#F5F5F5"
    }
    
    # Device-specific visual settings
    DESKTOP = {
        "figsize": (16, 9),
        "dpi": 300,
        "title_size": 20,
        "label_size": 15,
        "tick_size": 16,
        "legend_size": 14,
        "line_width": 4,
        "marker_size": 5,
        "grid": True,
        "tick_rotation": 90,
        "add_logo": True,
        "max_ticks": 25,
    }
    
    MOBILE = {
        "figsize": (8, 8),
        "dpi": 300,
        "title_size": 18,
        "label_size": 15,
        "tick_size": 15,
        "legend_size": 12,
        "line_width": 3,
        "marker_size": 5,
        "grid": True,
        "tick_rotation": 0,
        "max_ticks": 5,
        "add_logo": True,
    }
    
    def __init__(self, outdir: Path = Path("charts")):
        self.outdir = outdir
        self.outdir.mkdir(parents=True, exist_ok=True)
    
    def line_plot(self, metadata, stem: str, **kwargs):
        """
        Generate line charts for both desktop and mobile.
        
        Parameters:
        -----------
        metadata : dict
            Chart metadata (title, source, etc.)
        stem : str
            Base filename for outputs
        **kwargs : dict
            Keyword arguments passed to matplotlib plotting functions, including:
            - x: array-like - X-axis data points (required)
            - y: list of array-like - Y-axis data series (required if not using data/df)
            - data: DataFrame - Data source when using DataFrame column references
            - df: DataFrame - Alias for data
            - color/c: str or list - Colors for the plot series
            - linestyle/ls: str or list - Line styles for the plot series
            - linewidth/lw: float or list - Line widths for the plot series
            - marker: str or list - Marker styles for the plot series
            - markersize/ms: float or list - Marker sizes for the plot series
            - label: str or list - Labels for the legend
            - alpha: float or list - Transparency values for the plot series
            
        Returns:
        --------
        dict
            Dictionary containing the generated figure objects {'desktop': fig, 'mobile': fig}
        """
        # Generate desktop version
        desktop_fig = self._create_line_plot(
            metadata, 
            style=self.DESKTOP,
            **kwargs
        )
        self._save_chart(desktop_fig, f"{stem}_desktop", create_pptx=True)
        
        # Generate mobile version
        mobile_fig = self._create_line_plot(
            metadata, 
            style=self.MOBILE,
            **kwargs
        )
        self._save_chart(mobile_fig, f"{stem}_mobile", create_pptx=False)
        
        return {
            'desktop': desktop_fig,
            'mobile': mobile_fig
        }

    def _create_line_plot(self, metadata, style, **kwargs):
        """
        Internal method to create a line plot with appropriate styling.
        
        Parameters:
        -----------
        metadata : dict
            Chart metadata (title, source, etc.)
        style : dict
            Device-specific styling to apply
        **kwargs : dict
            Matplotlib-compatible plot parameters
            
        Returns:
        --------
        matplotlib.figure.Figure
            The generated figure
        """
        # Create figure and axis
        fig_kwargs = {
            'figsize': kwargs.pop('figsize', style["figsize"]),
            'dpi': kwargs.pop('dpi', style["dpi"]),
            'facecolor': kwargs.pop('facecolor', None),
            'edgecolor': kwargs.pop('edgecolor', None),
        }
        fig, ax = plt.subplots(**{k: v for k, v in fig_kwargs.items() if v is not None})
        
        # Handle different data input methods (matching matplotlib's flexibility)
        data = kwargs.pop('data', kwargs.pop('df', None))
        x = kwargs.pop('x', None)
        y = kwargs.pop('y', None)
        
        # Extract data directly from kwargs if not specified in x/y
        if x is None and 'x_data' in kwargs:
            x = kwargs.pop('x_data')
        if y is None and 'y_data_list' in kwargs:
            y = kwargs.pop('y_data_list')
        
        # Validate we have the necessary data
        if x is None:
            raise ValueError("X-axis data must be provided via 'x' or 'x_data'")
        
        # Handle DataFrame input
        if data is not None:
            if isinstance(x, str):
                x = data[x]
            if isinstance(y, (list, tuple)) and all(isinstance(item, str) for item in y):
                y = [data[col] for col in y]
            elif isinstance(y, str):
                y = [data[y]]
        
        # If y is not a list but x is, interpret x as y and generate a range for x
        if y is None and isinstance(x, (list, tuple, np.ndarray)) and not isinstance(x[0], (list, tuple, np.ndarray)):
            y = [x]
            x = np.arange(len(x))
        
        # If y is not provided or not a list, convert to a list for consistent handling
        if y is None:
            raise ValueError("Y-axis data must be provided via 'y', 'y_data_list', or as the first positional argument")
        elif not isinstance(y, (list, tuple)) or (len(y) > 0 and not hasattr(y[0], '__len__')):
            y = [y]
        
        # Process style parameters
        # Handle matplotlib aliases
        color = kwargs.pop('color', kwargs.pop('c', None))
        linestyle = kwargs.pop('linestyle', kwargs.pop('ls', None))
        linewidth = kwargs.pop('linewidth', kwargs.pop('lw', style["line_width"]))
        marker = kwargs.pop('marker', None)
        markersize = kwargs.pop('markersize', kwargs.pop('ms', style["marker_size"]))
        alpha = kwargs.pop('alpha', None)
        label = kwargs.pop('label', None)
        
        # Handle list or scalar inputs for style parameters
        colors = self._ensure_list(color, len(y), default=[self.COLORS["blue"], self.COLORS["light_blue"], 
                                                        self.COLORS["purple"], self.COLORS["orange"]])
        linestyles = self._ensure_list(linestyle, len(y), default=['-', '--', '-.', ':'])
        linewidths = self._ensure_list(linewidth, len(y))
        markers = self._ensure_list(marker, len(y), default=[None] * len(y))
        markersizes = self._ensure_list(markersize, len(y))
        alphas = self._ensure_list(alpha, len(y), default=[1.0] * len(y))
        
        # Handle labels
        if label is not None:
            labels = self._ensure_list(label, len(y))
        else:
            labels = kwargs.pop('labels', [f"Series {i+1}" for i in range(len(y))])
            labels = self._ensure_list(labels, len(y))
        
        # Plot each data series
        lines = []
        for i, y_data in enumerate(y):
            # Build plot kwargs for this series
            plot_kwargs = {
                'color': colors[i],
                'linestyle': linestyles[i],
                'linewidth': linewidths[i],
                'marker': markers[i],
                'markersize': markersizes[i],
                'alpha': alphas[i],
                'label': labels[i],
            }
            
            # Add any series-specific parameters from kwargs
            series_key = f"series_{i}"
            if series_key in kwargs:
                plot_kwargs.update(kwargs.pop(series_key))
            
            # Plot the series
            line, = ax.plot(x, y_data, **plot_kwargs)
            lines.append(line)
        
        # Apply title from metadata
        title = metadata.get('title', '')
        if title:
            ax.set_title(title, fontweight='bold', fontsize=style["title_size"])
        
        # Apply axes labels if provided
        xlabel = kwargs.pop('xlabel', None)
        if xlabel:
            ax.set_xlabel(xlabel, fontsize=style["label_size"])
        
        ylabel = kwargs.pop('ylabel', None)
        if ylabel:
            ax.set_ylabel(ylabel, fontsize=style["label_size"])
        
        # Apply grid
        grid = kwargs.pop('grid', style["grid"])
        ax.grid(grid)
        
        # Apply tick formatting
        tick_rotation = kwargs.pop('tick_rotation', style["tick_rotation"])
        tick_size = kwargs.pop('tick_size', style["tick_size"])
        
        plt.setp(ax.get_xticklabels(), rotation=tick_rotation, fontsize=tick_size)
        plt.setp(ax.get_yticklabels(), fontsize=tick_size)
        
        # Apply scale formatter if specified
        scale = kwargs.pop('scale', None)
        if scale:
            axis = kwargs.pop('axis_scale', 'y')
            self._apply_scale_formatter(ax, scale, axis)
        
        # Apply axis limits if provided
        if 'xlim' in kwargs:
            ax.set_xlim(kwargs.pop('xlim'))
        if 'ylim' in kwargs:
            ax.set_ylim(kwargs.pop('ylim'))
        
        # Apply custom ticks if specified
        if 'xticks' in kwargs:
            ticks = kwargs.pop('xticks')
            ax.set_xticks(ticks)
            # If custom tick labels are provided, use them
            if 'xticklabels' in kwargs:
                ax.set_xticklabels(kwargs.pop('xticklabels'))
            else:
                # Format as integers if they are whole numbers
                if all(float(x).is_integer() for x in ticks):
                    ax.set_xticklabels([f"{int(x)}" for x in ticks])
        
        # Set tick locators
        max_xticks = kwargs.pop('max_xticks', style.get("max_ticks", None))
        if max_xticks:
            ax.xaxis.set_major_locator(plt.MaxNLocator(max_xticks))
        
        # Add legend if any label is non-None and not empty
        if any(label for label in labels if label):
            legend_kwargs = {
                'fontsize': style["legend_size"],
                'loc': kwargs.pop('legend_loc', 'best'),
                'frameon': kwargs.pop('legend_frameon', True),
            }
            
            # Extract legend parameters if provided
            legend_params = kwargs.pop('legend_params', {})
            legend_kwargs.update(legend_params)
            
            # Create the legend
            legend = ax.legend(**legend_kwargs)
        
        # Apply any remaining kwargs to the axis using matplotlib's set methods
        for key, value in kwargs.items():
            try:
                # First try if it's a direct attribute
                if hasattr(ax, key):
                    setattr(ax, key, value)
                # Then try as a setter method
                else:
                    setter = getattr(ax, f"set_{key}", None)
                    if setter and callable(setter):
                        setter(value)
            except Exception as e:
                print(f"Warning: Could not set axis parameter '{key}': {e}")
        
        # Add footer elements (line, source, logo)
        self._add_footer(fig, metadata, style)
        
        # Apply tight layout, avoiding the footer area
        fig.tight_layout(rect=[0, style.get("footer_height", 0.1), 1, 1])
        
        return fig

    def _ensure_list(self, value, length, default=None):
        """
        Ensure a value is a list of the specified length.
        
        Parameters:
        -----------
        value : any
            The value to convert to a list
        length : int
            The desired length of the list
        default : list, optional
            Default list to use if value is None
            
        Returns:
        --------
        list
            A list of the specified length
        """
        if value is None:
            if default is None:
                return [None] * length
            return default
        
        if not isinstance(value, (list, tuple)):
            return [value] * length
        
        # Ensure the list is long enough by cycling
        result = list(value)
        while len(result) < length:
            result.extend(value[:length - len(result)])
        
        return result
        
    def bar_chart(self, x_data, y_data_list, metadata, stem: str) -> None:
        """Generate bar charts for both desktop and mobile."""
        # Similar implementation to line_plot but for bar charts
        # ...
        
    def waffle_chart(self, metadata, stem: str, **kwargs) -> None:
        """
        Generate waffle charts for both desktop and mobile.
        
        Parameters:
        -----------
        metadata : dict
            Chart metadata (title, source, etc)
        stem : str
            Base filename for outputs
        **kwargs : dict
            Keyword arguments for the waffle chart, including:
            - values: dict - Dictionary with labels as keys and values as values (required)
            - rows: int - Number of rows in the waffle chart (required)
            - colors: list - List of colors to use for different categories
            - legend: bool/dict - Whether to show a legend and legend parameters
            - And any other valid Waffle parameters
            
        Returns:
        --------
        dict
            Dictionary containing the generated figure objects {'desktop': fig, 'mobile': fig}
        """
        
        # Generate desktop version
        desktop_fig = self._create_waffle_chart(
            metadata, 
            style=self.DESKTOP,
            **kwargs
        )
        self._save_chart(desktop_fig, f"{stem}_desktop", create_pptx=True)
        
        # Generate mobile version
        mobile_fig = self._create_waffle_chart(
            metadata, 
            style=self.MOBILE,
            **kwargs
        )
        self._save_chart(mobile_fig, f"{stem}_mobile", create_pptx=False)
        
        return {
            'desktop': desktop_fig,
            'mobile': mobile_fig
        }

    def _create_waffle_chart(self, metadata, style, **kwargs):
        """
        Internal method to create a waffle chart with appropriate styling.
        
        Parameters:
        -----------
        metadata : dict
            Chart metadata (title, source, etc)
        style : dict
            Device-specific styling to apply
        **kwargs : dict
            Any parameters to pass to the Waffle constructor
            
        Returns:
        --------
        matplotlib.figure.Figure
            The generated figure
        """
        # Extract parameters that need special handling
        special_params = {}
        for param in ['legend', 'title_fontproperties', 'labels_fontproperties']:
            if param in kwargs:
                special_params[param] = kwargs.pop(param)
        
        # Start with style defaults relevant to Waffle
        waffle_defaults = {
            'figsize': style["figsize"]
        }
        
        # Create parameters by merging defaults with provided kwargs
        # (kwargs take precedence)
        waffle_params = {**waffle_defaults, **kwargs}
        
        # Create the waffle chart
        fig = plt.figure(
            FigureClass=Waffle,
            **waffle_params
        )
        
        # Apply title from metadata
        title = metadata.get('title', '')
        if title:
            fig.suptitle(
                title,
                fontsize=style["title_size"],
                fontweight='bold',
                y=0.98
            )
        
        # Process special parameters that need post-figure creation handling
        if 'legend' in special_params and special_params['legend']:
            legend_params = {}
            if isinstance(special_params['legend'], dict):
                legend_params = special_params['legend']
            
            # Apply the legend with styling
            legend = plt.legend(
                fontsize=legend_params.get('fontsize', style["legend_size"]),
                loc=legend_params.get('loc', 'best'),
                frameon=legend_params.get('frameon', True)
            )
            
            # Apply any additional legend styling
            for key, value in legend_params.items():
                if key not in ['fontsize', 'loc', 'frameon']:
                    try:
                        # Try to set the attribute if available
                        setter = getattr(legend, f"set_{key}", None)
                        if setter and callable(setter):
                            setter(value)
                    except Exception as e:
                        print(f"Warning: Could not set legend parameter '{key}': {e}")
        
        # Apply custom font properties if provided
        if 'title_fontproperties' in special_params:
            for title_obj in fig.findobj(plt.text.Text):
                if title_obj.get_text() == title:
                    title_obj.set_fontproperties(special_params['title_fontproperties'])
        
        # Add footer elements (line, source, logo)
        self._add_footer(fig, metadata, style)
        
        return fig
    
    def _save_chart(self, fig, filename, create_pptx=False):
        """Save chart as SVG, PNG, and optionally PPTX."""
        svg_path = self.outdir / f"{filename}.svg"
        png_path = self.outdir / f"{filename}.png"
        
        fig.savefig(svg_path, format="svg", dpi=300)
        fig.savefig(png_path, format="png", dpi=300)
        print(f"✓ saved {svg_path.name} and {png_path.name}")
        
        if create_pptx:
            pptx_path = self.outdir / f"{filename}.pptx"
            self._create_pptx(png_path, pptx_path)
            print(f"✓ saved {pptx_path.name}")
            
        plt.close(fig)
    
    def _apply_scale_formatter(self, ax, scale='billions', axis='y', decimals=0, prefix='$'):
        """Apply scale formatting to axis."""
        scales = {
            'billions': {'factor': 1e9, 'suffix': 'B'},
            'millions': {'factor': 1e6, 'suffix': 'M'},
            'thousands': {'factor': 1e3, 'suffix': 'K'},
            'percentage': {'factor': 0.01, 'suffix': '%', 'prefix': ''}
        }
        
        if scale in scales:
            scale_info = scales[scale]
            factor = scale_info['factor']
            suffix = scale_info.get('suffix', '')
            prefix = scale_info.get('prefix', prefix)
            
            def formatter(x, pos):
                return f'{prefix}{x/factor:.{decimals}f}{suffix}'
                
            from matplotlib.ticker import FuncFormatter
            if axis in ('y', 'both'):
                ax.yaxis.set_major_formatter(FuncFormatter(formatter))
            if axis in ('x', 'both'):
                ax.xaxis.set_major_formatter(FuncFormatter(formatter))
    
    def _add_footer(self, fig, metadata, style, bottom_margin=0.1):
        """
        Add footer elements to the figure: horizontal line, source text, and logo.
        Coordinates the placement of all footer elements.
        
        Parameters:
        -----------
        fig : matplotlib.figure.Figure
            The figure to add the footer to
        metadata : dict
            Dictionary containing chart metadata
        style : dict
            Style dictionary with visual settings
        bottom_margin : float, optional
            Bottom margin to reserve for the footer
        """
        # Check if footer should be displayed
        if metadata.get('footer', True) == False:
            return
        
        # Reserve space at the bottom for footer
        fig.subplots_adjust(bottom=bottom_margin)
        
        # Add horizontal spacer line
        spacer_y = bottom_margin / 2  # Place line halfway in the margin
        self._add_horizontal_spacer(fig, y_position=spacer_y, linewidth=1)
        
        # Add source if provided
        source_text = metadata.get('source')
        if source_text:
            self._add_source(fig, source_text)

        self._add_logo(fig)
    
    
    def _add_horizontal_spacer(self, fig, y_position=None, color=None, linewidth=0.5, extent=(0.02, 0.98)):
        """Add a horizontal line spacer to the figure."""
        # Set default values if not provided
        if y_position is None:
            y_position = 0.06
        
        if color is None:
            color = self.COLORS["medium_gray"]
        
        # Add the horizontal line
        fig.add_artist(plt.Line2D(
            [extent[0], extent[1]],  # x-positions (left, right)
            [y_position, y_position],  # y-positions (same for horizontal line)
            transform=fig.transFigure,  # Use figure coordinates
            color=color,
            linestyle='-',
            linewidth=linewidth
        ))
    
    def _add_logo(self, fig):
        """Add The Planetary Society logo to the figure."""
        try:
            
            logo_path = Path(__file__).parent.parent / "img" / "TPS_Logo_1Line-Black.png"
            if not logo_path.exists():
                return
                
            logo = mpimg.imread(str(logo_path))
            
            # Apply color mask if the logo has an alpha channel (RGBA)
            if logo.shape[2] == 4:  # RGBA format
                # Extract alpha channel
                alpha = logo[:, :, 3]
                
                # Create a color mask for the logo to better
                # match the chart colors
                hex_color = self.COLORS["lunar_dust"]
                rgb_color = np.array([
                    int(hex_color[1:3], 16) / 255.0,
                    int(hex_color[3:5], 16) / 255.0,
                    int(hex_color[5:7], 16) / 255.0
                ])
                
                # Create a new image where all non-transparent pixels are the specified color
                new_logo = np.zeros((logo.shape[0], logo.shape[1], 4))
                for i in range(3):  # RGB channels
                    new_logo[:, :, i] = rgb_color[i]
                new_logo[:, :, 3] = alpha  # Preserve original alpha
                
                logo = new_logo
            
            imagebox = OffsetImage(logo, zoom=0.08)
                        
            ab = AnnotationBbox(
                imagebox, 
                xy=(0.99, 0),  # Position at right, bottom corner
                xycoords='figure fraction',
                box_alignment=(1, 0),  # Align the right edge of the logo with the xy point
                frameon=False,
                pad=0  # No padding
            )
            fig.add_artist(ab)
            
            # Ensure the figure size accommodates the logo
            # This is important to prevent the logo from extending beyond the visible area
            fig.tight_layout(rect=[0, 0.09, 1, 1])
        except Exception as e:
            print(f"Warning: Could not add logo: {e}")
    
    def _add_source(self, fig, source_text):
        """Add source text to the bottom left of the figure."""
        if not source_text:
            return
            
        # Add text at the bottom left
        fig.text(
            0.02,  # x position (left side)
            0.01,  # y position (bottom)
            f"Source: {source_text}".upper(),
            fontsize=11,
            color=self.COLORS["lunar_dust"],
            ha='left',
            va='bottom'
        )
    
    def _create_pptx(self, png_path, pptx_path):
        """Create a PowerPoint file with the chart."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation()
        
        # Set slide size to 16x9
        prs.slide_width = Inches(13.33)  # 16 inches wide
        prs.slide_height = Inches(7.5)  # 9 inches tall
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.shapes.add_picture(str(png_path), Inches(0.5), Inches(0.5), width=Inches(12.33))
        prs.save(pptx_path)