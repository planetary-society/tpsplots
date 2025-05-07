"""tpsplots.base – shared chart infrastructure"""
from __future__ import annotations
from abc import ABC, abstractmethod
import os
from pathlib import Path
from typing import Tuple, Dict
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter
from pptx import Presentation
from pptx.util import Inches

from styles import ChartStyleFactory

# ------------------------------------------------------------------#
# Load TPS chart styles
plt.style.use(Path(__file__).parent / "style" / "tps.mplstyle")

class BaseChart(ABC):
    """Parent for all TPS charts: common style + export helpers."""

    def __init__(self, data_source=None, outdir: str | os.PathLike = "charts"):
        self.data_source = data_source
        self.outdir = Path(outdir)


    def apply_scale_formatter(self, ax, scale='billions', axis='y', decimals=1, prefix='$', suffix=None):
        """
        Format axis values with a scale factor, custom prefix and suffix.
        
        Args:
            ax: The matplotlib axis to format
            scale: Scale to use ('billions', 'millions', 'thousands', 'percentage', or numeric scale factor)
            axis: Which axis to format ('x', 'y', or 'both')
            decimals: Number of decimal places to show
            prefix: String prefix (e.g. '$')
            suffix: String suffix override (if None, uses default for scale)
        """
        # Set up scale-specific parameters
        scale_info = {
            'billions': {'factor': 1e9, 'default_suffix': 'B'},
            'millions': {'factor': 1e6, 'default_suffix': 'M'},
            'thousands': {'factor': 1e3, 'default_suffix': 'K'},
            'percentage': {'factor': 0.01, 'default_suffix': '%', 'prefix': ''}
        }
        
        # Get scale parameters or use custom scale factor
        if isinstance(scale, str) and scale.lower() in scale_info:
            scale_data = scale_info[scale.lower()]
            factor = scale_data['factor']
            # Use provided suffix or default for this scale
            suffix = suffix if suffix is not None else scale_data['default_suffix']
            # For percentage, override prefix if not explicitly provided
            if scale.lower() == 'percentage' and prefix == '$':
                prefix = scale_data.get('prefix', '')
        else:
            # Assume scale is a numeric scale factor
            try:
                factor = float(scale)
                suffix = suffix if suffix is not None else ''
            except (ValueError, TypeError):
                raise ValueError(f"Unknown scale: {scale}. Use 'billions', 'millions', 'thousands', "
                                f"'percentage', or a numeric scale factor.")
        
        # Create the formatter function
        def formatter(x, pos):
            if scale.lower() == 'percentage':
                # For percentage, we multiply by 100
                return f'{x*100:.{decimals}f}{suffix}'
            else:
                # For other scales, we divide by the factor
                return f'{prefix}{x/factor:.{decimals}f}{suffix}'
        
        # Apply the formatter to the specified axes
        fmt = FuncFormatter(formatter)
        
        if axis.lower() in ('y', 'both'):
            ax.yaxis.set_major_formatter(fmt)
        
        if axis.lower() in ('x', 'both'):
            ax.xaxis.set_major_formatter(fmt)

    # Convenience methods that call the unified formatter
    def apply_billions_formatter(self, ax, axis='y', decimals=1, prefix='$', suffix='B'):
        """Format axis values in billions."""
        return self.apply_scale_formatter(ax, 'billions', axis, decimals, prefix, suffix)

    def apply_millions_formatter(self, ax, axis='y', decimals=1, prefix='$', suffix='M'):
        """Format axis values in millions."""
        return self.apply_scale_formatter(ax, 'millions', axis, decimals, prefix, suffix)

    def apply_thousands_formatter(self, ax, axis='y', decimals=1, prefix='$', suffix='K'):
        """Format axis values in thousands."""
        return self.apply_scale_formatter(ax, 'thousands', axis, decimals, prefix, suffix)

    def apply_percentage_formatter(self, ax, axis='y', decimals=1, prefix='', suffix='%'):
        """Format axis values as percentages."""
        return self.apply_scale_formatter(ax, 'percentage', axis, decimals, prefix, suffix)

    # -------------------- exporting ---------------------------------

    def _export(
        self,
        fig: plt.Figure,
        stem: str,
        styles: Tuple[str, ...] = ("desktop", "mobile"),
        pptx: bool = True,
    ) -> None:
        """Save *fig* as SVG+PNG with each specified style; embed the desktop PNG in PPTX."""
        self.outdir.mkdir(parents=True, exist_ok=True)

        # Store original figure properties
        original_figsize = fig.get_size_inches()
        
        for style_name in styles:
            # Get style from factory
            style = ChartStyleFactory.make(style_name)
            print(style.stylesheets)
            # Apply the stylesheet directly
            plt.style.use(style.stylesheets)  # Reset to default before applying new style
            
            # Set figure size from style
            fig.set_size_inches(*style.figsize, forward=True)
            
            # Apply style-specific tweaks that can't be done via stylesheets
            if style_name == "mobile":
                for ax in fig.get_axes():
                    ax.xaxis.set_major_locator(plt.MaxNLocator(5))
                    ax.yaxis.set_major_locator(plt.MaxNLocator(5))
                    plt.setp(ax.get_xticklabels(), rotation=0)
            elif style_name == "desktop":
                for ax in fig.get_axes():
                    plt.setp(ax.get_xticklabels(), rotation=45, ha="right")
            
            # Add padding for logo if needed
            if style.add_logo and style.logo_path and style.logo_path.exists():
                fig.subplots_adjust(bottom=0.15)
            
            fig.tight_layout()
            
            # Add logo if specified
            if style.add_logo and style.logo_path and style.logo_path.exists():
                self._add_logo(fig, style.logo_path)
            
            # Save files
            svg = self.outdir / f"{stem}_{style_name}.svg"
            png = self.outdir / f"{stem}_{style_name}.png"
            fig.savefig(svg, format="svg", bbox_inches="tight", dpi=style.dpi)
            fig.savefig(png, format="png", bbox_inches="tight", dpi=style.dpi)
            print(f"✓ saved {svg.name} and {png.name}")

            # Create PPTX for desktop style
            if pptx and style_name == "desktop":
                self._png_to_pptx(png, self.outdir / f"{stem}.pptx")
            
            # Remove logo after saving if it was added
            if style.add_logo:
                for artist in fig.artists[:]:
                    if isinstance(artist, plt.matplotlib.offsetbox.AnnotationBbox):
                        artist.remove()
        
        # Restore original figure size
        fig.set_size_inches(*original_figsize, forward=True)
        plt.close(fig)

    def _add_logo(self, fig: plt.Figure, logo_path: Path) -> None:
        """Add The Planetary Society logo to the figure."""
        from matplotlib.offsetbox import OffsetImage, AnnotationBbox
        import matplotlib.image as mpimg
        
        # Load the logo
        logo = mpimg.imread(str(logo_path))
        
        # Create an OffsetImage with the logo
        imagebox = OffsetImage(logo, zoom=0.1)  # Adjust zoom as needed
        
        # Add the logo to the figure (bottom right corner)
        ab = AnnotationBbox(
            imagebox, 
            xy=(0.99, 0.01),  # Position (right, bottom)
            xycoords='figure fraction',
            box_alignment=(1, 0),  # Alignment (right, bottom)
            frameon=False
        )
        fig.add_artist(ab)

    @staticmethod
    def _png_to_pptx(png_path: Path, pptx_path: Path) -> None:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.shapes.add_picture(str(png_path), Inches(0.5), Inches(0.5), width=Inches(9))
        prs.save(pptx_path)