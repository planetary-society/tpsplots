import matplotlib.pyplot as plt
import matplotlib.font_manager as mt
from matplotlib.ticker import FormatStrFormatter
import pandas as pd
import os
from cycler import cycler
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns
import sys
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '../data_sources')))
from nasa_budget_data_source import Historical

# --- Base Chart Class ---
class BaseChart:
    """
    Base class for creating charts with a common theme and output methods.
    """
    
    CHARTS_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../charts"))

    # Planetary Society colors
    GREYS = [
        {"name": "Slushy Brine", "hex": "#F5F5F5", "rgb": (245, 245, 245)},
        {"name": "Comet Dust", "hex": "#C3C3C3", "rgb": (195, 195, 195)},
        {"name": "Lunar Soil", "hex": "#8C8C8C", "rgb": (140, 140, 140)},
        {"name": "Crater Shadow", "hex": "#414141", "rgb": (65, 65, 65)},
    ]
    
    PURPLES = [
        {"name": "Light Plasma", "hex": "#D8CDE1"},
        {"name": "Medium Plasma", "hex": "#B19BC3"},
        {"name": "Plasma Purple", "hex": "#643788"}
    ]
    
    BLUES = [
        {"name": "Neptune Blue", "hex": "#037CC2"},
        {"name": "Medium Neptune", "hex": "#80BDE0"},
        {"name": "Light Neptune", "hex": "#BFDEF0"},
    ]
    
    REDS = {
        "name": "Rocket Flame", "hex": "#FF5D47", "rgb": (255, 93, 71)
    }
    
    # Cycler colors for the theme
    PLANETARY_SOCIETY_COLORS = (
        '#037CC2', # Neptune Blue (Primary)
        '#643788', # Plasma Purple (Secondary)
        '#FF5D47', # Rocket Flame
        '#80BDE0', # Medium Neptune
        '#B19BC3', # Medium Plasma
        '#414141' # Crater Shadow
    )

    # Define the base theme structure
    BASE_THEME = {
        "name": "TPS",
        "description": "A theme based on The Planetary Society's brand guidelines.",
        "author": "Casey Dreier", # Assuming the original author, update if needed
        "theme": {
            "font.family": "Poppins",
            "font.sans-serif": ["Poppins"],
            "ytick.minor.visible": False,
            "xtick.minor.visible": False,
            "xtick.major.size": 0,
            "ytick.major.size": 0,
            "ytick.major.pad": 5,
            "xtick.major.pad": 5,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "yaxis.labellocation": "center",
            "xaxis.labellocation": "center",
            "axes.grid": True,
            "grid.linewidth": 0.5,
            "grid.alpha": 0.3,
            "ytick.labelsize": 9,
            "xtick.labelsize": 9,
            "grid.linestyle": "-",
            "axes.linewidth": 1.2,
            "grid.color": "#C3C3C3",
            "axes.edgecolor": "#C3C3C3",
            "axes.facecolor": "#F5F5F5",
            "figure.facecolor": "#F5F5F5",
            "axes.prop_cycle": cycler("color", PLANETARY_SOCIETY_COLORS),
            "font.size": 14,
            "axes.titlesize": 12,
        },
    }

    def __init__(self, data_source=None):
        """
        Initializes the BaseChart with optional data source.
        Applies the theme globally upon initialization.
        """
        self.data_source = data_source
        self._apply_theme()
        self.output_dir = "charts" # Default output directory

    def _apply_theme(self):
        """Applies the defined theme to matplotlib (which Seaborn uses)."""
        plt.style.use('default') # Reset to default before applying custom style
        plt.rcParams.update(self.BASE_THEME["theme"])
        # You can also use Seaborn's own styling if preferred,
        # but applying matplotlib rcParams first often works well.
        # sns.set_theme(style="whitegrid", rc=self.BASE_THEME["theme"]) # Example of Seaborn theme

    def _ensure_output_dir(self, chart_specific_dir=None):
        """Ensures the output directory exists."""
        if chart_specific_dir:
            full_path = os.path.join(self.output_dir, chart_specific_dir)
        else:
            full_path = self.output_dir

        if not os.path.exists(full_path):
            os.makedirs(full_path)
        return full_path

    def _save_chart(self, fig, filename, chart_specific_dir=None):
        """Saves the chart to SVG, PNG, and embeds PNG in a PPTX slide."""
        output_path = self._ensure_output_dir(chart_specific_dir)
        base_filepath = os.path.join(output_path, filename)

        # Save as SVG
        svg_filepath = f"{base_filepath}.svg"
        fig.savefig(svg_filepath, format='svg', bbox_inches='tight')
        print(f"Saved SVG: {svg_filepath}")

        # Save as PNG (300 dpi)
        png_filepath = f"{base_filepath}.png"
        fig.savefig(png_filepath, format='png', dpi=300, bbox_inches='tight')
        print(f"Saved PNG: {png_filepath}")

        # Embed PNG in PPTX
        pptx_filepath = f"{base_filepath}.pptx"
        self._embed_png_in_pptx(png_filepath, pptx_filepath, filename)
        print(f"Saved PPTX: {pptx_filepath}")

        plt.close(fig) # Close the figure to free memory

    def _embed_png_in_pptx(self, png_filepath, pptx_filepath, slide_title="Chart"):
        """Embeds a PNG image into a new PowerPoint slide."""
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6] # Use a blank slide layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add a title placeholder (optional, adjust position as needed)
        # For a blank layout, we might need to add a textbox manually
        # Or use a layout with a title placeholder if available and suitable

        # Add the image
        # Calculate position and size - adjust as needed for your desired layout
        # This is a basic placement - you might want to center it or size it differently
        img_path = png_filepath
        left = top = Inches(0.5)
        pic = slide.shapes.add_picture(img_path, left, top, width=Inches(9)) # Adjust width as needed

        # Save the presentation
        prs.save(pptx_filepath)

# --- Specific Chart Subclasses ---

class NASABudgetChart(BaseChart):
    """
    Charts specifically for NASA budget data.
    """
    def __init__(self):
        data_source = Historical() # Initialize the data source
        super().__init__(data_source=data_source)
        self.output_dir = os.path.join(self.CHARTS_DIR, "nasa_budget")

    def pbrs(self):
        """
        Generates a line chart of NASA PBR by Fiscal Year using Seaborn.
        """
        df = self.data_source.data()

        if df is None or df.empty:
            print("No data available to plot PBRs.")
            return

        # Create a matplotlib figure and axes, then pass the axes to seaborn
        fig, ax = plt.subplots(figsize=(10, 6)) # Adjust figure size as needed

        # Plotting the PBR data using Seaborn
        # Seaborn's lineplot automatically handles many styling aspects
        sns.lineplot(data=df, x="Fiscal Year", y="PBR", marker='o', ax=ax)

        # Formatting the plot (can still use matplotlib methods on the axes)
        ax.set_title("NASA PBR by Fiscal Year", fontsize=14, fontweight='bold')
        ax.set_xlabel("Fiscal Year", fontsize=10)
        ax.set_ylabel("PBR (Billions USD)", fontsize=10)

        # Improve x-axis ticks visibility if many years
        if len(df["Fiscal Year"]) > 15:
             ax.xaxis.set_major_locator(plt.MaxNLocator(15)) # Limit to max 15 ticks

        # Format y-axis labels as currency (optional, requires matplotlib.ticker)
        ax.yaxis.set_major_formatter(FormatStrFormatter('$%1.2fB'))

        plt.xticks(rotation=45, ha='right') # Rotate x-axis labels for better readability
        plt.tight_layout() # Adjust layout to prevent labels overlapping

        # Save the chart
        self._save_chart(fig, "nasa_pbr_by_year", chart_specific_dir="") # Use the subclass's output_dir

# --- Example Usage ---
if __name__ == "__main__":
    # Instantiate the chart class
    nasa_charts = NASABudgetChart()

    # Generate the PBR chart
    nasa_charts.pbrs()

    print("\nChart generation complete. Check the 'charts/nasa_budget' directory.")
