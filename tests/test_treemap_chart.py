"""Tests for the treemap chart configuration and renderer."""

from pathlib import Path
from zipfile import ZipFile

import pytest
from matplotlib import pyplot as plt
from matplotlib.patches import Rectangle
from PIL import Image
from pptx import Presentation
from pydantic import ValidationError

from tpsplots.models.chart_config import CHART_TYPES
from tpsplots.models.charts import CONFIG_REGISTRY
from tpsplots.models.charts.treemap import TreemapChartConfig
from tpsplots.views import VIEW_REGISTRY
from tpsplots.views.treemap_chart import TreemapChartView

VIKING_EXAMPLE = (
    Path(__file__).parent.parent / "yaml" / "examples" / "viking_cost_breakdown_treemap.yaml"
)
VIKING_COMPONENTS = [
    "Project Management",
    "Orbiter Spacecraft",
    "Orbiter Science",
    "Lander Spacecraft",
    "Lander Science",
    "Launch Vehicle",
    "Operations",
]


def test_treemap_is_registered_across_public_dispatch_tables():
    """Treemap config and view dispatch must be available through every registry."""
    assert CHART_TYPES["treemap"] == "treemap_plot"
    assert CONFIG_REGISTRY["treemap"].__name__ == "TreemapChartConfig"
    assert VIEW_REGISTRY["treemap_plot"].__name__ == "TreemapChartView"


def test_treemap_has_scaffold_template_and_generated_docs_example():
    from tpsplots.docs_generator import MINIMAL_EXAMPLES
    from tpsplots.templates import TEMPLATES

    assert "type: treemap" in TEMPLATES["treemap"]
    assert 'colors: ["Neptune Blue", "Plasma Purple", "Rocket Flame"]' in TEMPLATES["treemap"]
    assert 'edgecolor: "Polar White"' in TEMPLATES["treemap"]
    assert "type: treemap" in MINIMAL_EXAMPLES["treemap"]


def test_treemap_config_exposes_responsive_label_and_tile_defaults():
    config = TreemapChartConfig(
        output="costs",
        title="Project costs",
        labels="{{labels}}",
        values="{{values}}",
    )

    assert config.colors is None
    assert config.edgecolor == "Polar White"
    assert config.linewidth == 2.0
    assert config.alpha == 1.0
    assert config.show_labels is True
    assert config.show_values is False
    assert config.show_percentages is True
    assert config.value_format == "float"
    assert config.label_min_area_pct == 1.0
    assert config.label_wrap_length is None
    assert config.label_fontsize is None


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("linewidth", -0.1),
        ("alpha", -0.1),
        ("alpha", 1.1),
        ("label_min_area_pct", -0.1),
        ("label_min_area_pct", 100.1),
        ("label_wrap_length", 0),
        ("label_fontsize", 0),
    ],
)
def test_treemap_config_rejects_out_of_range_style_values(field, value):
    with pytest.raises(ValidationError):
        TreemapChartConfig(
            output="costs",
            title="Project costs",
            labels=["A"],
            values=[1],
            **{field: value},
        )


@pytest.mark.parametrize("geometry_key", ["x", "y", "xy", "width", "height"])
def test_treemap_config_rejects_geometry_in_matplotlib_escape_hatch(geometry_key):
    with pytest.raises(ValidationError, match="geometry"):
        TreemapChartConfig(
            output="costs",
            title="Project costs",
            labels=["A"],
            values=[1],
            matplotlib_config={geometry_key: 10},
        )


def test_normalize_items_stably_sorts_and_preserves_explicit_color_associations(tmp_path):
    view = TreemapChartView(outdir=tmp_path)

    items = view._normalize_items(
        labels=["first equal", "largest", "second equal"],
        values=[2, 5, 2],
        colors=["#111111", "#222222", "#333333"],
    )

    assert [(item.label, item.value, item.color) for item in items] == [
        ("largest", 5.0, "#222222"),
        ("first equal", 2.0, "#111111"),
        ("second equal", 2.0, "#333333"),
    ]


def test_normalize_items_assigns_default_cycle_after_value_sort(tmp_path):
    view = TreemapChartView(outdir=tmp_path)

    items = view._normalize_items(labels=["small", "large"], values=[1, 4])

    assert [item.label for item in items] == ["large", "small"]
    assert [item.color for item in items] == view._get_cycled_colors(2)


@pytest.mark.parametrize(
    ("labels", "values", "colors"),
    [
        ("A", [1], None),
        ({"A": 1}, [1], None),
        ([], [], None),
        (["A"], "1", None),
        (["A"], [1, 2], None),
        ([None], [1], None),
        (["A"], [True], None),
        (["A"], ["not numeric"], None),
        (["A"], [float("nan")], None),
        (["A"], [float("inf")], None),
        (["A"], [0], None),
        (["A"], [-1], None),
        (["A"], [1], []),
    ],
)
def test_normalize_items_rejects_invalid_flat_treemap_payloads(tmp_path, labels, values, colors):
    view = TreemapChartView(outdir=tmp_path)

    with pytest.raises(ValueError):
        view._normalize_items(labels=labels, values=values, colors=colors)


def test_invalid_payload_is_rejected_before_figure_allocation(tmp_path, monkeypatch):
    view = TreemapChartView(outdir=tmp_path)

    def fail_if_called(*_args, **_kwargs):
        raise AssertionError("figure allocation must happen after validation")

    monkeypatch.setattr(view, "_setup_figure", fail_if_called)
    with pytest.raises(ValueError, match="greater than zero"):
        view.create_figure(
            metadata={"title": "Invalid"}, labels=["A"], values=[0], device="desktop"
        )


def test_overflowing_value_total_is_rejected_before_figure_allocation(tmp_path, monkeypatch):
    view = TreemapChartView(outdir=tmp_path)

    def fail_if_called(*_args, **_kwargs):
        raise AssertionError("figure allocation must happen after validation")

    monkeypatch.setattr(view, "_setup_figure", fail_if_called)
    with pytest.raises(ValueError, match="finite total"):
        view.create_figure(
            metadata={"title": "Invalid"},
            labels=["A", "B"],
            values=[1e308, 1e308],
            device="desktop",
        )


def test_layout_items_fills_requested_aspect_ratio_with_proportional_areas(tmp_path):
    view = TreemapChartView(outdir=tmp_path)
    items = view._normalize_items(labels=["medium", "small", "large"], values=[3, 1, 6])

    tiles = view._layout_items(items, width=200.0, height=100.0)

    assert tiles[0][0].label == "large"
    assert tiles[0][1]["x"] == pytest.approx(0.0)
    assert tiles[0][1]["y"] == pytest.approx(0.0)
    assert [rect["dx"] * rect["dy"] for _, rect in tiles] == pytest.approx(
        [12000.0, 6000.0, 2000.0]
    )
    assert sum(rect["dx"] * rect["dy"] for _, rect in tiles) == pytest.approx(20000.0)
    assert all(
        0 <= rect["x"] <= 200
        and 0 <= rect["y"] <= 100
        and rect["x"] + rect["dx"] <= 200 + 1e-9
        and rect["y"] + rect["dy"] <= 100 + 1e-9
        for _, rect in tiles
    )


def test_format_tile_text_independently_controls_label_value_and_percentage(tmp_path):
    view = TreemapChartView(outdir=tmp_path)
    item = view._normalize_items(["Lander Spacecraft"], [451_700_000])[0]

    assert (
        view._format_tile_text(item, 1_057_900_000, True, True, True, "monetary", 10)
        == "Lander\nSpacecraft\n$452M\n42.7%"
    )
    assert (
        view._format_tile_text(item, 1_057_900_000, True, False, False, "monetary", 10)
        == "Lander\nSpacecraft"
    )
    assert (
        view._format_tile_text(item, 1_057_900_000, False, True, False, "monetary", 10) == "$452M"
    )
    assert (
        view._format_tile_text(item, 1_057_900_000, False, False, True, "monetary", 10) == "42.7%"
    )
    assert view._format_tile_text(item, 1_057_900_000, False, False, False, "monetary", 10) == ""


def test_contrast_text_color_uses_composited_tile_color(tmp_path):
    view = TreemapChartView(outdir=tmp_path)

    assert view._contrast_text_color("#643788", 1.0, "#F5F5F5") == "#FFFFFF"
    assert view._contrast_text_color("#037CC2", 1.0, "#F5F5F5") == "#000000"
    assert view._contrast_text_color("#BFDEF0", 1.0, "#F5F5F5") == "#000000"
    assert view._contrast_text_color("#000000", 0.1, "#FFFFFF") == "#000000"


def test_create_figure_draws_proportional_tiles_in_final_desktop_axes(tmp_path):
    view = TreemapChartView(outdir=tmp_path)

    fig = view.create_figure(
        metadata={"title": "Project costs", "source": "Test data"},
        device="desktop",
        labels=["small", "large", "medium"],
        values=[1, 6, 3],
        show_labels=False,
        show_percentages=False,
    )

    try:
        ax = fig.axes[0]
        rectangles = [patch for patch in ax.patches if isinstance(patch, Rectangle)]
        assert len(rectangles) == 3
        assert tuple(fig.get_size_inches()) == pytest.approx((16, 10))
        assert fig.dpi == pytest.approx(300)
        assert ax.yaxis_inverted()
        assert [patch.get_width() * patch.get_height() for patch in rectangles] == pytest.approx(
            [
                rectangles[0].get_width() * rectangles[0].get_height() * ratio
                for ratio in (1, 0.5, 1 / 6)
            ]
        )
    finally:
        plt.close(fig)


@pytest.mark.parametrize("device", ["desktop", "mobile", "social"])
def test_rendered_labels_are_suppressed_unless_they_fit_their_device_tile(tmp_path, device):
    view = TreemapChartView(outdir=tmp_path)
    fig = view.create_figure(
        metadata={"title": "Project costs", "source": "Test data"},
        device=device,
        labels=["Largest category", "Small category that cannot be labeled"],
        values=[99, 1],
        label_min_area_pct=2,
    )

    try:
        fig.canvas.draw()
        ax = fig.axes[0]
        visible_texts = [text for text in ax.texts if text.get_visible()]
        expected = "Largest\ncategory\n99.0%" if device == "mobile" else "Largest category\n99.0%"
        assert [text.get_text() for text in visible_texts] == [expected]

        renderer = fig.canvas.get_renderer()
        label_bbox = visible_texts[0].get_window_extent(renderer)
        patch = ax.patches[0]
        corners = ax.transData.transform(
            [
                (patch.get_x(), patch.get_y()),
                (patch.get_x() + patch.get_width(), patch.get_y() + patch.get_height()),
            ]
        )
        inset = 4 * fig.dpi / 72
        assert label_bbox.x0 >= min(corners[:, 0]) + inset
        assert label_bbox.x1 <= max(corners[:, 0]) - inset
        assert label_bbox.y0 >= min(corners[:, 1]) + inset
        assert label_bbox.y1 <= max(corners[:, 1]) - inset
    finally:
        plt.close(fig)


def test_flattened_matplotlib_escape_hatch_overrides_typed_tile_style(tmp_path):
    view = TreemapChartView(outdir=tmp_path)
    fig = view._create_chart(
        metadata={"title": "Project costs"},
        style=view.DESKTOP,
        labels=["A"],
        values=[1],
        colors=["#111111"],
        show_labels=False,
        show_percentages=False,
        facecolor="#ABCDEF",
        zorder=7,
    )

    try:
        patch = fig.axes[0].patches[0]
        assert patch.get_facecolor()[:3] == pytest.approx((171 / 255, 205 / 255, 239 / 255))
        assert patch.get_zorder() == 7
    finally:
        plt.close(fig)


@pytest.mark.parametrize(
    ("device", "expected_size", "expected_dpi"),
    [
        ("desktop", (16, 10), 300),
        ("mobile", (8, 9), 300),
        ("social", (8, 4.2), 150),
    ],
)
def test_treemap_geometry_tracks_final_axes_aspect_ratio(
    tmp_path, device, expected_size, expected_dpi
):
    view = TreemapChartView(outdir=tmp_path)
    fig = view.create_figure(
        metadata={"title": "Responsive treemap", "source": "Test data"},
        device=device,
        labels=["A", "B", "C"],
        values=[6, 3, 1],
        show_labels=False,
        show_percentages=False,
    )

    try:
        fig.canvas.draw()
        ax = fig.axes[0]
        bbox = ax.get_window_extent(fig.canvas.get_renderer())
        xlim = ax.get_xlim()
        ylim = ax.get_ylim()
        data_ratio = abs((xlim[1] - xlim[0]) / (ylim[1] - ylim[0]))
        assert tuple(fig.get_size_inches()) == pytest.approx(expected_size)
        assert fig.dpi == pytest.approx(expected_dpi)
        assert data_ratio == pytest.approx(bbox.width / bbox.height)
        assert all(
            patch.get_x() >= 0
            and patch.get_y() >= 0
            and patch.get_x() + patch.get_width() <= max(xlim) + 1e-9
            and patch.get_y() + patch.get_height() <= max(ylim) + 1e-9
            for patch in ax.patches
        )
    finally:
        plt.close(fig)


def test_viking_example_resolves_seven_nonoverlapping_nominal_cost_components():
    from tpsplots.processors.yaml_chart_processor import YAMLChartProcessor

    processor = YAMLChartProcessor(VIKING_EXAMPLE)

    assert processor.config.chart.labels == VIKING_COMPONENTS
    assert processor.config.chart.values == pytest.approx(
        [48.2, 193.3, 23.9, 451.7, 158.2, 78.6, 104.0]
    )
    assert sum(processor.config.chart.values) == pytest.approx(1057.9)
    assert list(processor.data["metadata"]["column_sums"]) == VIKING_COMPONENTS


def test_viking_example_generates_all_responsive_outputs(tmp_path):
    from tpsplots.processors.yaml_chart_processor import YAMLChartProcessor

    processor = YAMLChartProcessor(VIKING_EXAMPLE, outdir=tmp_path)
    result = processor.generate_chart()

    assert "files" in result
    assert {path.name for path in tmp_path.iterdir()} == {
        "viking_cost_breakdown_treemap_desktop.svg",
        "viking_cost_breakdown_treemap_desktop.png",
        "viking_cost_breakdown_treemap.pptx",
        "viking_cost_breakdown_treemap_mobile.svg",
        "viking_cost_breakdown_treemap_mobile.png",
        "viking_cost_breakdown_treemap_social.png",
    }

    # generate_chart no longer returns the device figures (they are saved and
    # closed internally). Re-render each device to inspect the drawn labels.
    from tpsplots.processors.render_pipeline import build_render_context

    ctx = build_render_context(processor.config, processor.data, log_conflicts=False)
    view = VIEW_REGISTRY[ctx.chart_type_v1](outdir=tmp_path)
    expected_largest_labels = {"Lander Spacecraft", "Orbiter Spacecraft", "Lander Science"}
    for device in ("desktop", "mobile", "social"):
        fig = view.create_figure(
            metadata=ctx.resolved_metadata, device=device, **ctx.resolved_params
        )
        try:
            fig.canvas.draw()
            visible_labels = {
                text.get_text().replace("\n", " ").rsplit(" ", 1)[0]
                for text in fig.axes[0].texts
                if text.get_visible()
            }
            assert expected_largest_labels.issubset(visible_labels)
        finally:
            plt.close(fig)

    desktop_png = tmp_path / "viking_cost_breakdown_treemap_desktop.png"
    mobile_png = tmp_path / "viking_cost_breakdown_treemap_mobile.png"
    social_png = tmp_path / "viking_cost_breakdown_treemap_social.png"
    pptx_path = tmp_path / "viking_cost_breakdown_treemap.pptx"
    for image_path, expected_size in (
        (desktop_png, (4800, 3000)),
        (mobile_png, (2400, 2700)),
        (social_png, (1200, 630)),
    ):
        with Image.open(image_path) as image:
            assert image.size == expected_size

    with ZipFile(pptx_path) as archive:
        embedded_images = [name for name in archive.namelist() if name.startswith("ppt/media/")]
        assert len(embedded_images) == 1
        assert archive.read(embedded_images[0]) == desktop_png.read_bytes()

    presentation = Presentation(pptx_path)
    assert presentation.slide_width / presentation.slide_height == pytest.approx(16 / 9, rel=5e-4)
    picture = presentation.slides[0].shapes[0]
    assert picture.height == presentation.slide_height
    assert picture.left == pytest.approx((presentation.slide_width - picture.width) / 2, abs=1)
